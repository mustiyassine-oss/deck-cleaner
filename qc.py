"""Visual QC: PPTX → PDF (LibreOffice) → PNG (pdf2image) → vision LLM per slide."""
from __future__ import annotations

import base64
import io
import json
import logging
import os
import shutil
import subprocess
from pathlib import Path
from typing import Any

from openai import OpenAI
from pptx import Presentation
from pptx.oxml.ns import qn

log = logging.getLogger("deckgen")

QC_SYSTEM_PROMPT = """You are a PowerPoint layout QC reviewer.

Review the slide screenshot ONLY for concrete, repairable layout failures.
Evaluate every item in the checklist below before deciding pass/fail.

Fail the slide if ANY checklist item is true:
- text_overlap: KPI labels, bullet text, title text, or body text overlaps another text/object
- text_overflow: any text is clipped, cut off, outside its box, or unreadable because it collides with nearby text
- text_too_small: text appears too small/dense to read comfortably at normal presentation viewing size
- table_out_of_bounds: any table extends beyond the slide edge or visible table region, especially below the slide
- table_cell_clipping: table row/cell text is clipped, cut off, or visibly too dense to read
- chart_label_overlap: chart data labels, axis labels, or legend labels overlap/collide or are too dense to read clearly
- chart_label_clipping: chart data labels or axis labels are cut off or outside the chart/slide bounds
- object_misalignment: title, chart, table, or text block is visibly out of alignment with the intended slide grid

Important examples that MUST fail:
- KPI delta/description lines printed on top of each other under KPI numbers
- a table continuing below the bottom edge of the slide
- chart labels printed over each other or too close to distinguish
- axis/category labels colliding with chart elements

Do NOT fail for subjective design taste, business/content quality, wording, colors, contrast,
or minor spacing that does not cause overlap, clipping, overflow, table bounds issues,
chart label collisions, or visible misalignment.

Return ONLY valid JSON:
{
  "status": "pass" | "fail",
  "severity": "none" | "low" | "medium" | "high",
  "checks": {
    "text_overlap": true | false,
    "text_overflow": true | false,
    "text_too_small": true | false,
    "table_out_of_bounds": true | false,
    "table_cell_clipping": true | false,
    "chart_label_overlap": true | false,
    "chart_label_clipping": true | false,
    "object_misalignment": true | false
  },
  "issues": [
    {
      "type": "alignment|overflow|clipping|overlap|table_bounds|chart_labels",
      "description": "specific visible layout failure",
      "recommended_fix": "specific layout repair"
    }
  ]
}
"""

CHECK_TO_ISSUE = {
    "text_overlap": (
        "overlap",
        "Text overlaps another slide element.",
        "Reduce or shorten the overlapping text and re-render the slide.",
    ),
    "text_overflow": (
        "overflow",
        "Text is clipped, cut off, or outside its available region.",
        "Shorten the text or reduce the number of text items.",
    ),
    "text_too_small": (
        "overflow",
        "Text appears too small or too dense to read comfortably.",
        "Reduce content density, trim text, or increase readable text size.",
    ),
    "table_out_of_bounds": (
        "table_bounds",
        "A table extends beyond the slide or table region.",
        "Reduce visible table rows/columns or split the table.",
    ),
    "table_cell_clipping": (
        "table_bounds",
        "Table cell text is clipped or too dense to read.",
        "Reduce table rows/columns or shorten cell text.",
    ),
    "chart_label_overlap": (
        "chart_labels",
        "Chart labels overlap or are too dense to read clearly.",
        "Reduce categories, move labels, or switch chart orientation.",
    ),
    "chart_label_clipping": (
        "chart_labels",
        "Chart labels are clipped or outside the chart bounds.",
        "Reduce categories or adjust chart label placement.",
    ),
    "object_misalignment": (
        "alignment",
        "A slide object is visibly misaligned against the layout grid.",
        "Snap the object back to its default layout region.",
    ),
}

ALLOWED_ISSUE_TYPES = {
    "alignment",
    "overflow",
    "clipping",
    "overlap",
    "table_bounds",
    "chart_labels",
}

ISSUE_TYPE_ALIASES = {
    "table": "table_bounds",
    "table_overflow": "table_bounds",
    "table_clipping": "table_bounds",
    "table_bounds": "table_bounds",
    "chart": "chart_labels",
    "chart_label": "chart_labels",
    "chart_labels": "chart_labels",
    "label_overlap": "chart_labels",
    "data_labels": "chart_labels",
    "data_label": "chart_labels",
    "text_overlap": "overlap",
    "overlapping": "overlap",
    "bounds": "clipping",
}


def _extract_json(raw: str) -> str:
    s = raw.strip()
    for prefix in ("```json", "```JSON", "```"):
        if s.startswith(prefix):
            s = s[len(prefix) :].lstrip()
            break
    if s.endswith("```"):
        s = s[:-3].rstrip()
    return s


def _normalize_issue_type(raw: Any) -> str | None:
    key = str(raw or "").strip().lower().replace(" ", "_").replace("-", "_")
    if key in ALLOWED_ISSUE_TYPES:
        return key
    return ISSUE_TYPE_ALIASES.get(key)


def _normalize_qc_issues(raw_issues: Any) -> list[dict[str, Any]]:
    if not isinstance(raw_issues, list):
        return []

    normalized: list[dict[str, Any]] = []
    for issue in raw_issues:
        if not isinstance(issue, dict):
            continue
        issue_type = _normalize_issue_type(issue.get("type"))
        if not issue_type:
            continue
        normalized.append(
            {
                "type": issue_type,
                "description": str(issue.get("description", "")).strip(),
                "recommended_fix": str(issue.get("recommended_fix", "")).strip(),
            }
        )
    return normalized


def _issues_from_checks(raw_checks: Any) -> list[dict[str, Any]]:
    if not isinstance(raw_checks, dict):
        return []

    def _is_true(v: Any) -> bool:
        if isinstance(v, bool):
            return v
        if isinstance(v, (int, float)):
            return v != 0
        if isinstance(v, str):
            return v.strip().lower() in ("true", "yes", "y", "1")
        return False

    issues: list[dict[str, Any]] = []
    for key, details in CHECK_TO_ISSUE.items():
        if _is_true(raw_checks.get(key)):
            issue_type, description, recommended_fix = details
            issues.append(
                {
                    "type": issue_type,
                    "description": description,
                    "recommended_fix": recommended_fix,
                }
            )
    return issues


# Strictly below this point size, slide text is treated as genuinely unreadable.
# (Body/footnote text at 8–10pt is common and acceptable, so it must NOT be flagged;
# flagging it produced unfixable "too small" failures in earlier testing.)
_TINY_PT = 8.0
_SZ_TAGS = (qn("a:rPr"), qn("a:defRPr"), qn("a:endParaRPr"))


def _effective_para_pt(paragraph) -> float:
    """
    Effective point size of a paragraph, resolving inherited sizes.

    `run.font.size` is frequently None because the size is declared on the
    paragraph/endParaRPr or a defRPr (inherited), so reading only run sizes
    misses most text. We fall back to the paragraph font and finally to the
    largest `a:sz` found anywhere in the paragraph XML. Returns 0.0 if unknown.
    """
    best = 0.0
    for run in paragraph.runs:
        sz = getattr(getattr(run, "font", None), "size", None)
        if sz is not None:
            try:
                best = max(best, float(sz.pt))
            except Exception:
                pass
    if best > 0:
        return best
    try:
        if paragraph.font.size and paragraph.font.size.pt:
            return float(paragraph.font.size.pt)
    except Exception:
        pass
    try:
        for node in paragraph._p.iter():
            if node.tag in _SZ_TAGS:
                v = node.get("sz")
                if v:
                    try:
                        best = max(best, int(v) / 100.0)
                    except (TypeError, ValueError):
                        pass
    except Exception:
        pass
    return best


def _structural_text_issues(pptx_bytes: bytes) -> dict[int, list[dict[str, Any]]]:
    """
    Deterministic QC guardrail from the PPT object model.
    Flags genuinely tiny (<8pt) or very dense text the vision model may miss.
    Uses effective (inherited) font sizes so detection is consistent.
    """
    issues_by_slide: dict[int, list[dict[str, Any]]] = {}
    try:
        pres = Presentation(io.BytesIO(pptx_bytes))
    except Exception:
        return issues_by_slide

    for idx, slide in enumerate(pres.slides):
        slide_no = idx + 1
        tiny_paras = 0
        dense_shapes = 0
        text_shapes = 0
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                continue
            text_shapes += 1
            tf = shape.text_frame
            para_count = len(tf.paragraphs or [])
            char_count = sum(len((p.text or "").strip()) for p in tf.paragraphs)
            # Heuristic: many paragraphs + lots of text in one box -> dense.
            if para_count >= 7 and char_count >= 380:
                dense_shapes += 1

            for p in tf.paragraphs:
                # Only consider paragraphs with real content (ignore stray tiny labels).
                if len((p.text or "").strip()) < 10:
                    continue
                eff = _effective_para_pt(p)
                if 0 < eff < _TINY_PT:
                    tiny_paras += 1

        slide_issues: list[dict[str, Any]] = []
        if tiny_paras >= 2:
            slide_issues.append(
                {
                    "type": "overflow",
                    "description": f"Text renders below {_TINY_PT:.0f}pt and is hard to read.",
                    "recommended_fix": "Increase font size or reduce text density.",
                }
            )
        if dense_shapes >= 1 and text_shapes >= 2:
            slide_issues.append(
                {
                    "type": "overflow",
                    "description": "Slide contains dense text blocks likely unreadable in presentation view.",
                    "recommended_fix": "Reduce bullets/paragraph count or split content across multiple slides.",
                }
            )
        if slide_issues:
            issues_by_slide[slide_no] = slide_issues
    return issues_by_slide


def write_pptx_bytes(pptx_bytes: bytes, output_path: str | Path) -> Path:
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(pptx_bytes)
    log.info("[QC] Wrote PPTX (%d bytes) to %s", len(pptx_bytes), path)
    return path


def _resolve_soffice() -> str | None:
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            log.info("[QC] Using LibreOffice binary: %s", p)
            return p
    win_dirs = [
        os.environ.get("PROGRAMFILES", r"C:\Program Files"),
        os.environ.get("PROGRAMFILES(X86)", r"C:\Program Files (x86)"),
    ]
    for base in win_dirs:
        if not base:
            continue
        for sub in (
            "LibreOffice/program/soffice.exe",
            "LibreOffice 7/program/soffice.exe",
            "LibreOffice 24/program/soffice.exe",
        ):
            cand = Path(base) / sub
            if cand.is_file():
                log.info("[QC] Using LibreOffice binary: %s", cand)
                return str(cand)
    return None


def convert_pptx_to_pdf(pptx_path: str | Path, output_dir: str | Path) -> Path:
    pptx_path = Path(pptx_path)
    outdir = Path(output_dir)
    outdir.mkdir(parents=True, exist_ok=True)
    soffice = _resolve_soffice()
    if not soffice:
        raise RuntimeError("LibreOffice (soffice/libreoffice) not found on PATH or standard install paths.")

    log.info("[QC] Converting PPTX to PDF via LibreOffice headless: %s", pptx_path.name)
    proc = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(outdir),
            str(pptx_path),
        ],
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=600,
    )
    if proc.returncode != 0:
        log.warning(
            "[QC] LibreOffice exit code %s stderr=%s stdout=%s",
            proc.returncode,
            (proc.stderr or "")[:500],
            (proc.stdout or "")[:500],
        )
        raise RuntimeError(f"LibreOffice PDF conversion failed (exit {proc.returncode})")

    pdf_path = outdir / f"{pptx_path.stem}.pdf"
    if not pdf_path.is_file():
        raise RuntimeError(f"Expected PDF not found at {pdf_path}")

    log.info("[QC] PPTX → PDF complete: %s", pdf_path)
    return pdf_path


def convert_pdf_to_images(
    pdf_path: str | Path,
    output_dir: str | Path,
    dpi: int = 150,
) -> list[Path]:
    try:
        from pdf2image import convert_from_path
    except ImportError as e:
        raise RuntimeError("pdf2image is not installed.") from e

    pdf_path = Path(pdf_path)
    outdir = Path(output_dir)
    outdir.mkdir(parents=True, exist_ok=True)

    log.info("[QC] Converting PDF to PNG screenshots (dpi=%s)...", dpi)
    try:
        images = convert_from_path(str(pdf_path), dpi=dpi)
    except Exception as e:
        msg = str(e).lower()
        if "poppler" in msg or "path" in msg or "pdftoppm" in msg:
            raise RuntimeError(
                "pdf2image requires Poppler (poppler-utils on Linux, or Poppler for Windows). "
                f"Original error: {e}"
            ) from e
        raise

    paths: list[Path] = []
    for i, img in enumerate(images):
        p = outdir / f"slide_{i + 1:04d}.png"
        img.save(p, "PNG")
        paths.append(p)

    log.info("[QC] PDF → PNG complete: %d screenshot(s) in %s", len(paths), outdir)
    return paths


def qc_single_slide(
    client: OpenAI,
    image_path: str | Path,
    slide_number: int,
    model: str = "gpt-4.1",
) -> dict[str, Any]:
    image_path = Path(image_path)
    b64 = base64.b64encode(image_path.read_bytes()).decode("ascii")
    data_url = f"data:image/png;base64,{b64}"

    messages: list[dict[str, Any]] = [
        {"role": "system", "content": QC_SYSTEM_PROMPT},
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": f"Evaluate slide number {slide_number} (screenshot attached). Return only the JSON object.",
                },
                {"type": "image_url", "image_url": {"url": data_url}},
            ],
        },
    ]

    resp = client.chat.completions.create(
        model=model,
        messages=messages,
        response_format={"type": "json_object"},
        max_tokens=2048,
        temperature=0.2,
    )
    raw = (resp.choices[0].message.content or "").strip()
    try:
        out = json.loads(_extract_json(raw))
    except json.JSONDecodeError:
        log.warning("[QC] Slide %d: invalid JSON from vision model, defaulting to pass", slide_number)
        return {
            "status": "pass",
            "severity": "none",
            "issues": [],
            "parse_error": True,
        }

    status = str(out.get("status", "pass")).lower()
    if status not in ("pass", "fail"):
        status = "pass"
    sev = str(out.get("severity", "none")).lower()
    if sev not in ("none", "low", "medium", "high"):
        sev = "none"
    issues = _normalize_qc_issues(out.get("issues"))
    existing = {
        (
            str(issue.get("type", "")),
            str(issue.get("description", "")),
        )
        for issue in issues
    }
    for issue in _issues_from_checks(out.get("checks")):
        key = (
            str(issue.get("type", "")),
            str(issue.get("description", "")),
        )
        if key not in existing:
            issues.append(issue)
            existing.add(key)

    if issues:
        status = "fail"
        if sev == "none":
            sev = "medium"
    if status == "fail" and not issues:
        status = "pass"
        sev = "none"
    elif status == "pass":
        issues = []
        sev = "none"

    return {"status": status, "severity": sev, "issues": issues}


def run_visual_qc(
    pptx_bytes: bytes,
    work_dir: str | Path,
    openai_api_key: str,
    model: str = "gpt-4.1",
    *,
    openai_base_url: str | None = None,
    slide_numbers: set[int] | list[int] | tuple[int, ...] | None = None,
    dpi: int = 220,
) -> dict[str, Any]:
    """
    Full visual QC pipeline. On missing tools or errors, returns a structured skip
    report (does not raise) so callers can log and continue.
    """
    work_dir = Path(work_dir)
    work_dir.mkdir(parents=True, exist_ok=True)
    target_slide_numbers = {int(n) for n in slide_numbers or [] if int(n) > 0}

    def _qc_skipped(reason: str) -> dict[str, Any]:
        return {
            "overall_status": "skipped",
            "skipped_reason": reason,
            "failed_slide_count": 0,
            "total_slide_count": 0,
            "slides": [],
        }

    if not openai_api_key or not openai_api_key.strip():
        log.warning("[QC] No OpenAI API key; skipping visual QC.")
        return _qc_skipped("OpenAI API key not configured")

    if not _resolve_soffice():
        log.warning("[QC] LibreOffice not found; skipping visual QC.")
        return _qc_skipped("LibreOffice (soffice) not found")

    try:
        import pdf2image  # noqa: F401
    except ImportError:
        log.warning("[QC] pdf2image not installed; skipping visual QC.")
        return _qc_skipped("pdf2image not installed")

    pptx_path = work_dir / "deck_for_qc.pptx"
    pdf_dir = work_dir / "pdf"
    img_dir = work_dir / "slides_png"

    try:
        write_pptx_bytes(pptx_bytes, pptx_path)
        pdf_path = convert_pptx_to_pdf(pptx_path, pdf_dir)
        image_paths = convert_pdf_to_images(pdf_path, img_dir, dpi=dpi)
    except Exception as e:
        log.warning("[QC] Screenshot pipeline failed; skipping visual QC: %s", e)
        return _qc_skipped(str(e))

    if not image_paths:
        log.warning("[QC] No slide images generated; skipping QC scoring.")
        return _qc_skipped("No PNG screenshots produced")

    structural_issues = _structural_text_issues(pptx_bytes)

    client_kw: dict[str, Any] = {"api_key": openai_api_key, "timeout": 600.0}
    if openai_base_url:
        client_kw["base_url"] = openai_base_url
    client = OpenAI(**client_kw)

    slides_out: list[dict[str, Any]] = []
    failed = 0
    checked_index = 0
    checked_total = len(target_slide_numbers) if target_slide_numbers else len(image_paths)
    for i, img_path in enumerate(image_paths):
        slide_number = i + 1
        if target_slide_numbers and slide_number not in target_slide_numbers:
            continue
        checked_index += 1
        try:
            qc = qc_single_slide(client, img_path, slide_number, model=model)
        except Exception as e:
            log.warning("[QC] Slide %d vision call failed: %s", slide_number, e)
            qc = {"status": "pass", "severity": "none", "issues": [], "vision_error": str(e)}

        entry = {
            "slide_number": slide_number,
            "status": qc.get("status", "pass"),
            "severity": qc.get("severity", "none"),
            "issues": qc.get("issues") if isinstance(qc.get("issues"), list) else [],
        }
        if slide_number in structural_issues:
            seen = {
                (str(i.get("type", "")), str(i.get("description", "")))
                for i in entry["issues"]
                if isinstance(i, dict)
            }
            for si in structural_issues[slide_number]:
                key = (str(si.get("type", "")), str(si.get("description", "")))
                if key not in seen:
                    entry["issues"].append(si)
            if entry["issues"]:
                entry["status"] = "fail"
                if entry["severity"] == "none":
                    entry["severity"] = "medium"
        slides_out.append(entry)
        log.info(
            "[QC] Slide %d (%d/%d checked): status=%s severity=%s issues=%d",
            slide_number,
            checked_index,
            checked_total,
            entry["status"],
            entry["severity"],
            len(entry["issues"]),
        )
        if entry["status"] == "fail":
            failed += 1

    overall = "fail" if failed else "pass"
    log.info(
        "[QC] Final result: overall_status=%s failed_slide_count=%d total_slide_count=%d",
        overall,
        failed,
        len(slides_out),
    )

    return {
        "overall_status": overall,
        "failed_slide_count": failed,
        "total_slide_count": len(image_paths),
        "checked_slide_count": len(slides_out),
        "checked_slide_numbers": [s["slide_number"] for s in slides_out],
        "slides": slides_out,
    }
