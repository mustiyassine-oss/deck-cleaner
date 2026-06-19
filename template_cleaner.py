"""Deterministic cleanup pass for existing PPTX templates."""
from __future__ import annotations

import io
import math
from dataclasses import dataclass
from functools import lru_cache
from typing import Any

from pptx import Presentation
from pptx.enum.chart import XL_DATA_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

try:
    from PIL import ImageFont
    _PIL_OK = True
except Exception:  # pragma: no cover - Pillow is a dependency, but stay safe
    _PIL_OK = False

_EMU_PER_PT = 12700


@dataclass(frozen=True)
class _Spacing:
    """
    Spacing/threshold values derived from the actual slide size so the cleaner
    adapts to any template/report dimensions instead of using hardcoded inches.
    All values are EMU. Ratios are calibrated against a standard 13.33x7.5" slide.
    """
    slide_w: int
    slide_h: int

    @property
    def gap(self) -> int:
        return max(1, int(self.slide_h * 0.0067))          # ~0.05"

    @property
    def min_y_overlap(self) -> int:
        return max(1, int(self.slide_h * 0.0053))          # ~0.04"

    @property
    def pad_v(self) -> int:
        return max(1, int(self.slide_h * 0.004))           # ~0.03"

    @property
    def pad_h(self) -> int:
        return max(1, int(self.slide_w * 0.0135))          # ~0.18"

    @property
    def slack(self) -> int:
        return max(1, int(self.slide_h * 0.0067))          # ~0.05"

    @property
    def below_tol(self) -> int:
        return max(1, int(self.slide_h * 0.008))           # ~0.06"

    @property
    def min_container_side(self) -> int:
        return max(1, int(self.slide_h * 0.04))            # ~0.3"

    @property
    def footer_margin(self) -> int:
        return max(1, int(self.slide_h * 0.08))            # ~0.6"

    @property
    def header_margin(self) -> int:
        # Protected top band (logo + title + subtitle). Text is never pulled up
        # into it, so de-overlap can't slide a title under the logo.
        return max(1, int(self.slide_h * 0.22))            # ~1.65"

    def max_shift(self, aggressive: bool) -> int:
        return int(self.slide_h * (0.3 if aggressive else 0.2))  # ~2.25" / ~1.5"


@lru_cache(maxsize=256)
def _load_font(family: str, size_px: int):
    """
    Best-effort TrueType font for width measurement. Tries the run's font family,
    then common sans-serif fallbacks. Returns None if nothing loads (callers then
    fall back to a character-count heuristic). Never raises.
    """
    if not _PIL_OK or size_px <= 0:
        return None
    candidates: list[str] = []
    fam = (family or "").strip()
    if fam:
        compact = fam.replace(" ", "")
        candidates += [fam, fam + ".ttf", compact + ".ttf", fam + ".otf", compact + ".otf"]
    candidates += ["arial.ttf", "Arial.ttf", "calibri.ttf", "DejaVuSans.ttf", "LiberationSans-Regular.ttf"]
    for cand in candidates:
        try:
            return ImageFont.truetype(cand, size_px)
        except Exception:
            continue
    return None


def _measure_text_width_emu(text: str, family: str, fs_pt: float) -> int | None:
    """Measure rendered text width in EMU using real font metrics, or None."""
    if not text:
        return 0
    size_px = int(round(fs_pt))
    font = _load_font(family or "", size_px)
    if font is None:
        return None
    try:
        # Font built at size=fs_pt px => 1px == 1pt, so width(px) == width(pt).
        width_pt = font.getlength(text)
        return int(width_pt * _EMU_PER_PT)
    except Exception:
        return None


def _issue_types_for_slide(qc_report: dict[str, Any], slide_number: int) -> list[str]:
    slides = qc_report.get("slides")
    if not isinstance(slides, list):
        return []
    for s in slides:
        if not isinstance(s, dict):
            continue
        if int(s.get("slide_number", -1)) != slide_number:
            continue
        issues = s.get("issues") if isinstance(s.get("issues"), list) else []
        out: list[str] = []
        for i in issues:
            if isinstance(i, dict):
                t = str(i.get("type", "")).strip().lower()
                if t:
                    out.append(t)
        return out
    return []


def _reduce_text_shape(shape, *, aggressive: bool) -> bool:
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return False
    changed = False
    tf = shape.text_frame
    try:
        if tf.word_wrap is not True:
            tf.word_wrap = True
            changed = True
        if tf.auto_size != MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            changed = True
        if tf.margin_top != Emu(20000):
            tf.margin_top = Emu(20000)
            changed = True
        if tf.margin_bottom != Emu(20000):
            tf.margin_bottom = Emu(20000)
            changed = True
        if tf.margin_left != Emu(30000):
            tf.margin_left = Emu(30000)
            changed = True
        if tf.margin_right != Emu(30000):
            tf.margin_right = Emu(30000)
            changed = True
    except Exception:
        pass

    step = 2 if aggressive else 1
    min_pt = 9 if aggressive else 10
    for p in tf.paragraphs:
        try:
            if p.space_before != Pt(0):
                p.space_before = Pt(0)
                changed = True
            if p.space_after != Pt(0):
                p.space_after = Pt(0)
                changed = True
            target_spacing = 1.0 if aggressive else 1.05
            if p.line_spacing != target_spacing:
                p.line_spacing = target_spacing
                changed = True
        except Exception:
            pass
        for run in p.runs:
            try:
                if run.font.size and run.font.size.pt:
                    new_pt = max(min_pt, run.font.size.pt - step)
                    if new_pt < run.font.size.pt:
                        run.font.size = Pt(new_pt)
                        changed = True
            except Exception:
                pass
    return changed


def _clean_table_shape(shape, *, aggressive: bool) -> bool:
    if not getattr(shape, "has_table", False):
        return False
    changed = False
    table = shape.table
    total_rows = len(table.rows)
    total_cols = len(table.columns)
    max_rows_fit = max(2, int(shape.height // Emu(260000)))
    keep_rows = min(total_rows, max_rows_fit if not aggressive else max_rows_fit - 1)
    keep_rows = max(2, keep_rows)

    for r in range(total_rows):
        for c in range(total_cols):
            cell = table.cell(r, c)
            tf = cell.text_frame
            if r >= keep_rows:
                if tf.text:
                    tf.clear()
                    changed = True
                continue
            try:
                if tf.word_wrap is not True:
                    tf.word_wrap = True
                    changed = True
                if tf.auto_size != MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    changed = True
                for p in tf.paragraphs:
                    for run in p.runs:
                        if run.font.size and run.font.size.pt:
                            target = 9 if aggressive else 10
                            if run.font.size.pt > target:
                                run.font.size = Pt(target)
                                changed = True
            except Exception:
                pass
    return changed


def _clean_chart_shape(shape, *, aggressive: bool) -> bool:
    if not getattr(shape, "has_chart", False):
        return False
    changed = False
    chart = shape.chart
    try:
        for series in chart.series:
            if aggressive and series.has_data_labels:
                series.has_data_labels = False
                changed = True
            if not aggressive and series.has_data_labels:
                dl = series.data_labels
                if dl.position != XL_DATA_LABEL_POSITION.OUTSIDE_END:
                    dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    changed = True
    except Exception:
        pass

    try:
        if chart.has_legend:
            if chart.legend.include_in_layout is not False:
                chart.legend.include_in_layout = False
                changed = True
    except Exception:
        pass

    try:
        cat_ax = chart.category_axis
        target = Pt(9 if aggressive else 10)
        if cat_ax.tick_labels.font.size != target:
            cat_ax.tick_labels.font.size = target
            changed = True
    except Exception:
        pass
    try:
        val_ax = chart.value_axis
        target = Pt(9 if aggressive else 10)
        if val_ax.tick_labels.font.size != target:
            val_ax.tick_labels.font.size = target
            changed = True
    except Exception:
        pass
    return changed


def _shape_box(shape) -> list[int] | None:
    """Return [left, top, width, height] in EMU, or None if not positionable."""
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
    except (TypeError, ValueError):
        return None
    if width <= 0 or height <= 0:
        return None
    return [left, top, width, height]


def _has_visible_text(shape) -> bool:
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return False
    try:
        return bool(shape.text_frame.text and shape.text_frame.text.strip())
    except Exception:
        return False


_SZ_TAGS = (qn("a:rPr"), qn("a:defRPr"), qn("a:endParaRPr"))


def _xml_font_pt(paragraph) -> float:
    """
    Largest `a:sz` (pt) declared anywhere in the paragraph XML.

    Critical: many templates set the font size on the paragraph's endParaRPr or
    a defRPr rather than on the run's rPr, so run.font.size returns None even
    though the text clearly renders at e.g. 44pt. We must read the XML directly.
    """
    best = 0.0
    try:
        for node in paragraph._p.iter():
            if node.tag in _SZ_TAGS:
                sz = node.get("sz")
                if sz:
                    try:
                        best = max(best, int(sz) / 100.0)
                    except (TypeError, ValueError):
                        pass
    except Exception:
        pass
    return best


def _paragraph_font_pt(paragraph) -> float:
    """Best-effort font size (pt) for a paragraph, falling back to a default."""
    size = 0.0
    for run in paragraph.runs:
        try:
            if run.font.size and run.font.size.pt:
                size = max(size, float(run.font.size.pt))
        except Exception:
            pass
    if size <= 0:
        try:
            if paragraph.font.size and paragraph.font.size.pt:
                size = float(paragraph.font.size.pt)
        except Exception:
            pass
    if size <= 0:
        size = _xml_font_pt(paragraph)
    return size if size > 0 else 18.0


def _paragraph_spacing_emu(paragraph) -> int:
    """Inter-paragraph spacing (space_before + space_after) in EMU, 0 if unset.

    Point-valued spacing (``a:spcPts``) is reported by python-pptx as an EMU
    Length; percentage spacing (``a:spcPct``) has no fixed EMU and is ignored.
    Templates routinely add a few points before/after each bullet, which adds up
    over a list and is a real part of the rendered text height."""
    total = 0
    for attr in ("space_before", "space_after"):
        try:
            v = getattr(paragraph, attr)
        except Exception:
            v = None
        if v is None:
            continue
        try:
            total += int(v)
        except (TypeError, ValueError):
            pass
    return total


def _paragraph_font_name(paragraph) -> str:
    """Best-effort font family for a paragraph (for width measurement)."""
    for run in paragraph.runs:
        try:
            if run.font.name:
                return str(run.font.name)
        except Exception:
            pass
    try:
        if paragraph.font.name:
            return str(paragraph.font.name)
    except Exception:
        pass
    return ""


def _segment_line_count(seg: str, family: str, fs_pt: float, usable_w: int) -> int:
    """How many wrapped lines a single (unbroken) text segment needs."""
    if not seg:
        return 1
    measured = _measure_text_width_emu(seg, family, fs_pt)
    if measured is not None and measured > 0:
        return max(1, math.ceil(measured / usable_w))
    # Fallback heuristic: average glyph advance ~0.55 * font-size.
    glyph_w = max(1.0, fs_pt * 0.55 * _EMU_PER_PT)
    chars_per_line = max(1, int(usable_w / glyph_w))
    return max(1, math.ceil(len(seg) / chars_per_line))


def _effective_text_height(shape, width_emu: int, declared_h: int) -> int:
    """
    Effective rendered text height (EMU) for collision detection.

    PowerPoint/LibreOffice render text that overflows a shape's declared box, so
    the declared `height` can't be trusted when a title wraps to extra lines.
    BUT a large single-line value (e.g. a 28pt KPI number in a short, centered
    box) does *not* really overflow — its glyph height fits even though a naive
    line-height estimate exceeds the box. So we only report overflow when the
    text needs MORE lines than the box can hold; otherwise we keep the declared
    height. This avoids falsely shoving KPI labels out of their tiles while still
    catching genuine multi-line overflow (forced breaks / long wraps).
    """
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return declared_h
    tf = shape.text_frame
    try:
        margin_l = int(tf.margin_left) if tf.margin_left is not None else 91440
        margin_r = int(tf.margin_right) if tf.margin_right is not None else 91440
        margin_t = int(tf.margin_top) if tf.margin_top is not None else 45720
        margin_b = int(tf.margin_bottom) if tf.margin_bottom is not None else 45720
    except Exception:
        margin_l = margin_r = 91440
        margin_t = margin_b = 45720

    usable_w = max(1, width_emu - margin_l - margin_r)
    content = float(margin_t + margin_b)
    total_lines = 0
    total_spacing = 0.0
    max_line_h = 0.0
    for p in tf.paragraphs:
        text = p.text or ""
        fs = _paragraph_font_pt(p)
        family = _paragraph_font_name(p)
        # Honor forced line breaks (<a:br> -> "\v", and "\n") as hard newlines,
        # then wrap each resulting segment. Templates frequently use a soft break
        # in titles like "SHORT-TERM\v(Next 30 Days)". Wrapping uses real font
        # metrics when available, else a character-count heuristic.
        segments = (text.replace("\v", "\n")).split("\n")
        lines = 0
        for seg in segments:
            lines += _segment_line_count(seg, family, fs, usable_w)
        lines = max(1, lines)
        try:
            ls = p.line_spacing
            ls = float(ls) if isinstance(ls, (int, float)) else 1.0
        except Exception:
            ls = 1.0
        line_h = fs * 1.2 * _EMU_PER_PT * (ls if ls and ls > 0 else 1.0)
        content += lines * line_h
        total_lines += lines
        total_spacing += _paragraph_spacing_emu(p)
        max_line_h = max(max_line_h, line_h)

    if max_line_h <= 0:
        return declared_h
    content += total_spacing
    # Capacity is the line count the box holds AFTER reserving inter-paragraph
    # spacing. With no spacing this is identical to the old `declared_h // line_h`
    # (so existing behaviour is unchanged); with real spcBef/spcAft it correctly
    # tightens the budget so a spaced-out bullet list is recognised as overflowing.
    usable_h = declared_h - total_spacing
    capacity_lines = max(1, int(usable_h // max_line_h)) if usable_h > 0 else 0
    if total_lines <= capacity_lines:
        # Text fits the lines the box was sized for -> no overflow. Use the
        # SMALLER of the declared box and the rendered ink: a generously-sized
        # box that the text underfills (e.g. a tall "key takeaways" panel only
        # half full) must not be treated as colliding with content sitting in its
        # empty lower half. ``content`` already includes the line-height padding,
        # so for a single big line (KPI value) it stays >= declared and this is a
        # no-op — preserving the existing KPI guard.
        return min(declared_h, int(content))
    return max(declared_h, int(content))


def _group_child_bottom(group_shape, fallback_bottom: int) -> int:
    """Bottom boundary (EMU) of a group's child coordinate space."""
    try:
        el = group_shape._element
        grp_pr = el.find(qn("p:grpSpPr"))
        xfrm = grp_pr.find(qn("a:xfrm")) if grp_pr is not None else None
        ch_off = xfrm.find(qn("a:chOff")) if xfrm is not None else None
        ch_ext = xfrm.find(qn("a:chExt")) if xfrm is not None else None
        if ch_off is not None and ch_ext is not None:
            return int(ch_off.get("y")) + int(ch_ext.get("cy"))
    except Exception:
        pass
    return fallback_bottom


def _slide_containers(slide, spacing: _Spacing) -> list[list[int]]:
    """
    Collect "container" shapes — the card/tile autoshapes that text sits inside
    (e.g. KPI rounded-rectangles, the recommendation cards). Excludes text boxes,
    thin accent bars, and the full-slide background. Returns [l, t, w, h] boxes.
    """
    slide_area = max(1, spacing.slide_w * spacing.slide_h)
    min_side = spacing.min_container_side
    out: list[list[int]] = []
    for shape in slide.shapes:
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
        except Exception:
            continue
        if _has_visible_text(shape):
            continue
        box = _shape_box(shape)
        if box is None:
            continue
        if box[2] < min_side or box[3] < min_side:
            continue
        if box[2] * box[3] > slide_area * 0.6:  # skip full-slide backgrounds
            continue
        out.append(box)
    return out


def _find_container(box: list[int], containers: list[list[int]], spacing: _Spacing) -> list[int] | None:
    """
    Smallest card/tile that a *small* text box belongs inside: it must enclose
    the box horizontally and the box must sit within the card vertically (or just
    below it). Big, near-full-height boxes (e.g. a card's bullet body) are not
    treated as "inside" a container — only small elements like KPI labels or card
    titles are, so we never clamp/move the main body text.
    """
    slack = spacing.slack
    below_tol = spacing.below_tol
    bl, bt, bw, bh = box
    br = bl + bw
    best = None
    best_area = None
    for c in containers:
        cl, ct, cw, ch = c
        cr, cb = cl + cw, ct + ch
        if cl > bl + slack or cr < br - slack:
            continue  # must enclose horizontally
        if bt < ct - slack or bt > cb + below_tol:
            continue  # box must start inside the card (or just below it)
        if bh > ch * 0.6:
            continue  # too tall to be a small in-card element (e.g. body text)
        area = cw * ch
        if best is None or area < best_area:
            best, best_area = c, area
    return best


# ── Structural keep-out regions (side panels / header-footer bands) ──────────
#
# The de-overlap pass only understands text-vs-text. But content also collides
# with *decorative structural rectangles*: a dark sidebar (e.g. a "Key Takeaways"
# panel) or the footer/title band. Those are solid auto-shapes/pictures, not text
# and not the slide background, so nothing flagged them before. These helpers
# detect such regions and measure how far a content shape pokes into one from the
# outside, so both the detector (pdf_geometry) and the fixer can act on them.


def _keepout_regions(slide, spacing: _Spacing) -> list[tuple[str, list[int]]]:
    """
    Structural keep-out rectangles: tall side PANELS and wide BANDS at top/bottom.
    Solid decorative auto-shapes/pictures only — not text, not the full-slide
    background. Returns [(kind, [l, t, w, h])] with kind in
    {"panel_left", "panel_right", "band_top", "band_bottom"}.
    """
    sw, sh = spacing.slide_w, spacing.slide_h
    slide_area = max(1, sw * sh)
    out: list[tuple[str, list[int]]] = []
    for shape in slide.shapes:
        # A structural wall/band is a solid rectangle or a picture. Decorative
        # ovals/diamonds/arcs are NOT panels — treating a big background circle as
        # a side panel previously shoved a whole column (title included) sideways
        # under the logo, which the verify gate then reverted wholesale.
        if not _is_card_background(shape):
            continue
        if _has_visible_text(shape):
            continue  # a panel that holds its own text is a content card, not a wall
        box = _shape_box(shape)
        if box is None:
            continue
        l, t, w, h = box
        area = w * h
        if area > slide_area * 0.85 or area < slide_area * 0.05:
            continue  # full-slide background / too small to be structural
        cx = l + w / 2.0
        tall = h >= sh * 0.55 and w <= sw * 0.6
        wide = w >= sw * 0.55 and h <= sh * 0.22
        # A real band hugs a slide edge. Requiring this prevents wide *content*
        # rows in the middle of a slide (e.g. stacked message cards) from being
        # mistaken for a header/footer band — which previously set a bogus footer
        # line and let the footer fixer shrink unrelated shapes.
        hugs_top = t <= sh * 0.05 and (t + h) <= sh * 0.22
        hugs_bottom = (t + h) >= sh * 0.95 and t >= sh * 0.78
        if tall and cx >= sw * 0.5:
            out.append(("panel_right", box))
        elif tall and cx < sw * 0.5:
            out.append(("panel_left", box))
        elif wide and hugs_bottom:
            out.append(("band_bottom", box))
        elif wide and hugs_top:
            out.append(("band_top", box))
    return out


def _intrusion_amount(kind: str, region: list[int], box: list[int], tol: int) -> int:
    """
    EMU by which `box` pokes into `region` *from the outside* (0 if it doesn't).

    "From the outside" means the box's center is on the content side of the wall
    but an edge crosses it — so content legitimately sitting ON the panel (center
    inside) is never flagged, only stuff spilling into it.
    """
    rl, rt, rw, rh = region
    rr, rb = rl + rw, rt + rh
    l, t, w, h = box
    r, b = l + w, t + h
    cx, cy = l + w / 2.0, t + h / 2.0
    if kind in ("panel_right", "panel_left"):
        if min(b, rb) - max(t, rt) <= 0:
            return 0  # no vertical overlap with the panel
        if kind == "panel_right":
            return max(0, r - rl) if (cx < rl and r > rl + tol) else 0
        return max(0, rr - l) if (cx > rr and l < rr - tol) else 0
    if min(r, rr) - max(l, rl) <= 0:
        return 0  # no horizontal overlap with the band
    if kind == "band_bottom":
        return max(0, b - rt) if (cy < rt and b > rt + tol) else 0
    return max(0, rb - t) if (cy > rb and t < rb - tol) else 0


def _intruder_boxes(slide, spacing: _Spacing, region_boxes: list[list[int]]) -> list[tuple[Any, list[int]]]:
    """Content shapes that could spill into a keep-out region: text boxes and
    card-sized auto-shapes (never the structural panels/bands themselves)."""
    slide_area = max(1, spacing.slide_w * spacing.slide_h)
    out: list[tuple[Any, list[int]]] = []
    for shape in slide.shapes:
        box = _shape_box(shape)
        if box is None or box in region_boxes:
            continue
        # Full-width thin rules (decorative divider lines spanning the slide) are
        # not content spilling into anything — ignore them as intruders.
        if box[2] >= spacing.slide_w * 0.9 and box[3] <= spacing.slide_h * 0.04:
            continue
        if _has_visible_text(shape):
            out.append((shape, box))
            continue
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
        except Exception:
            continue
        if box[2] * box[3] > slide_area * 0.25:
            continue  # too big to be a content card
        out.append((shape, box))
    return out


def _count_keepout_intrusions(slide, spacing: _Spacing) -> int:
    """Number of distinct content shapes spilling into any keep-out region."""
    regions = _keepout_regions(slide, spacing)
    if not regions:
        return 0
    region_boxes = [r[1] for r in regions]
    intruders = _intruder_boxes(slide, spacing, region_boxes)
    threshold = spacing.gap
    # Exempt only the actual title strip (the detected top band), not the whole
    # protected header zone — otherwise body content like a KPI row that starts
    # just under the title would be wrongly exempted from side-panel checks.
    header_cut = max(
        (b[1] + b[3] for k, b in regions if k == "band_top"),
        default=int(spacing.slide_h * 0.10),
    )
    hit: set[int] = set()
    for kind, region in regions:
        is_side = kind in ("panel_left", "panel_right")
        for shape, box in intruders:
            # A wide title in the header strip overlaps a same-coloured side panel
            # but its glyphs don't — don't treat header-band content as spilling
            # onto a side panel.
            if is_side and (box[1] + box[3] / 2.0) <= header_cut:
                continue
            if _intrusion_amount(kind, region, box, threshold) > threshold:
                hit.add(id(shape))
    return len(hit)


def _resolve_container(shapes, bound_bottom: int, spacing: _Spacing, *, aggressive: bool,
                       containers: list[list[int]] | None = None) -> bool:
    """
    Resolve vertical overlaps among shapes that share one coordinate space
    (a slide, or a single group's child space).

    For each overlapping pair the lower box is pushed down using estimated
    effective text height. When that downward travel is clamped (by the box's
    container card/tile or the slide bottom) and overlap remains, the *upper* box
    is pulled up by the leftover amount — bounded by its own card top or the slide
    top. This bidirectional move resolves overlaps inside fixed-height cards that a
    one-directional push could not separate. Footers/page numbers are never moved.
    """
    gap = spacing.gap
    max_shift = spacing.max_shift(aggressive)
    min_overlap_ratio = 0.45
    min_y_overlap = spacing.min_y_overlap          # ignore hairline/false overlaps
    pad_v = spacing.pad_v                           # padding inside a card
    footer_zone_top = bound_bottom - spacing.footer_margin  # leave footers/page numbers alone
    header_zone_bottom = spacing.header_margin      # protected top band (logo/title/subtitle)

    items: list[list[Any]] = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
        except Exception:
            pass
        if not _has_visible_text(shape):
            continue
        box = _shape_box(shape)
        if box is None:
            continue
        eff_h = _effective_text_height(shape, box[2], box[3])
        cont = _find_container(box, containers, spacing) if containers else None
        # [shape, [l, t, w, declared_h], effective_h, abs_shift, container, orig_top]
        items.append([shape, box, eff_h, 0, cont, box[1]])

    if len(items) < 2:
        return False

    for _ in range(8):
        items.sort(key=lambda it: (it[1][1], it[1][0]))  # by top, then left
        moved = False
        n = len(items)
        for i in range(n):
            a_box, a_eff = items[i][1], items[i][2]
            for j in range(i + 1, n):
                b_box, b_eff = items[j][1], items[j][2]
                # horizontal overlap (same column / card)
                x_overlap = min(a_box[0] + a_box[2], b_box[0] + b_box[2]) - max(a_box[0], b_box[0])
                if x_overlap <= 0:
                    continue
                if x_overlap < min(a_box[2], b_box[2]) * min_overlap_ratio:
                    continue
                # vertical overlap using EFFECTIVE (rendered) heights
                y_overlap = min(a_box[1] + a_eff, b_box[1] + b_eff) - max(a_box[1], b_box[1])
                if y_overlap <= min_y_overlap:
                    continue

                desired = y_overlap + gap

                # ── 1) push the lower box (j) down as far as allowed ──
                # Footers / page numbers are never moved, but we must NOT skip the
                # pair: instead leave j put and pull the upper box (i) up below.
                j_is_footer = b_box[1] >= footer_zone_top
                achieved_down = 0
                if not j_is_footer:
                    allowed_down = desired
                    if items[j][3] + allowed_down > max_shift:
                        allowed_down = max_shift - items[j][3]
                    new_top_j = b_box[1] + max(0, allowed_down)
                    cont_j = items[j][4]
                    if cont_j is not None:
                        cont_max_top = cont_j[1] + cont_j[3] - pad_v - b_box[3]
                        if new_top_j > cont_max_top:
                            new_top_j = cont_max_top
                    if new_top_j + b_box[3] > bound_bottom:
                        new_top_j = bound_bottom - b_box[3]
                    # Never push a body box down into the footer band (this is what
                    # ejected a card title into the footer). Only clamp when it
                    # wouldn't shove the box back above its current position.
                    if new_top_j + b_box[3] > footer_zone_top:
                        clamped = footer_zone_top - b_box[3]
                        if clamped >= b_box[1]:
                            new_top_j = clamped
                    achieved_down = new_top_j - b_box[1]
                    if achieved_down > 0:
                        b_box[1] = new_top_j
                        items[j][3] += achieved_down
                        moved = True

                # ── 2) if overlap remains, pull the upper box (i) up ──
                if j_is_footer:
                    # Clear the upper box's bottom above the footer's TOP. The
                    # effective-overlap metric understates the needed lift when the
                    # footer is short, leaving a hairline residual, so target the
                    # footer's top directly.
                    residual = (a_box[1] + a_eff) - b_box[1] + gap
                else:
                    residual = desired - max(0, achieved_down)
                # Never move the header/title/logo band, and never pull body text up
                # into it (that is what slid a title under the logo before).
                if residual > min_y_overlap and header_zone_bottom <= a_box[1] < footer_zone_top:
                    allowed_up = residual
                    if items[i][3] + allowed_up > max_shift:
                        allowed_up = max_shift - items[i][3]
                    floor_top = header_zone_bottom
                    cont_i = items[i][4]
                    if cont_i is not None:
                        floor_top = max(floor_top, cont_i[1] + pad_v)
                    new_top_i = a_box[1] - max(0, allowed_up)
                    if new_top_i < floor_top:
                        new_top_i = floor_top
                    up_move = a_box[1] - new_top_i
                    if up_move > 0:
                        a_box[1] = new_top_i
                        items[i][3] += up_move
                        moved = True
        if not moved:
            break

    changed = False
    for shape, box, _eff, _shift, _cont, orig_top in items:
        if box[1] != orig_top:
            try:
                shape.top = Emu(box[1])
                changed = True
            except Exception:
                pass
    return changed


def _is_left_aligned(shape) -> bool:
    try:
        for p in shape.text_frame.paragraphs:
            if p.alignment in (PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.JUSTIFY):
                return False
    except Exception:
        pass
    return True


def _fit_text_in_containers(slide, containers: list[list[int]], spacing: _Spacing) -> bool:
    """
    Tidy text inside card/tile containers:
      * pull a box back up if it hangs below the card bottom, and
      * give left-aligned text a bit more left padding so it doesn't hug the
        card edge (skips centered/right text like KPI numbers).
    """
    if not containers:
        return False
    pad_v = spacing.pad_v
    pad_h = spacing.pad_h
    footer_zone_top = spacing.slide_h - spacing.footer_margin  # leave footers/page numbers alone
    changed = False
    for shape in slide.shapes:
        if not _has_visible_text(shape):
            continue
        box = _shape_box(shape)
        if box is None:
            continue
        if box[1] >= footer_zone_top:
            continue
        cont = _find_container(box, containers, spacing)
        if cont is None:
            continue
        cl, ct, cw, ch = cont
        cr, cb = cl + cw, ct + ch
        left, top, w, h = box

        # Vertical: if the box hangs below the card, pull it up to fit.
        if top + h > cb - pad_v:
            new_top = cb - pad_v - h
            if new_top < ct + pad_v:
                new_top = ct + pad_v
            if new_top < top:
                try:
                    shape.top = Emu(new_top)
                    top = new_top
                    changed = True
                except Exception:
                    pass

        # Horizontal: left-aligned text shouldn't hug the card's left edge.
        if _is_left_aligned(shape) and (left - cl) < pad_h:
            new_left = min(cl + pad_h, cr - w)  # keep box within the card
            if new_left > left:
                try:
                    shape.left = Emu(new_left)
                    changed = True
                except Exception:
                    pass
    return changed


def _balance_kpi_cards(slide, containers: list[list[int]], spacing: _Spacing) -> bool:
    """
    Give KPI tiles breathing room between the value and its caption.

    A KPI tile is a small card holding exactly two vertically-stacked text boxes:
    a large value on top (e.g. "173") and a smaller label below ("TOTAL ARTICLES").
    Templates often place the label hard against the number. This pass *only
    increases* the gap toward an adaptive target (scaled to the value's font size),
    splitting the extra space between nudging the value up and the label down while
    keeping both inside the card's padding. It never reduces existing spacing and
    never moves a box outside its card, so it cannot create overlaps — overlap
    resolution stays the de-overlap pass's job; this is pure intra-card padding.
    """
    if not containers:
        return False
    pad_v = spacing.pad_v
    slack = spacing.slack

    # Collect text boxes once, with a representative (largest) font size.
    text_items: list[tuple[Any, list[int], float]] = []
    for shape in slide.shapes:
        if not _has_visible_text(shape):
            continue
        box = _shape_box(shape)
        if box is None:
            continue
        pt = 0.0
        try:
            for p in shape.text_frame.paragraphs:
                if (p.text or "").strip():
                    pt = max(pt, _paragraph_font_pt(p))
        except Exception:
            pass
        text_items.append((shape, box, pt))

    changed = False
    for cl, ct, cw, ch in containers:
        cr, cb = cl + cw, ct + ch
        # Target small KPI tiles only (not full-width / tall recommendation cards).
        if ch > spacing.slide_h * 0.28 or cw > spacing.slide_w * 0.40:
            continue

        members = [
            [shape, list(box), pt]
            for shape, box, pt in text_items
            if box[0] >= cl - slack
            and box[0] + box[2] <= cr + slack
            and box[1] >= ct - slack
            and box[1] + box[3] <= cb + spacing.below_tol
            and box[3] <= ch * 0.7
        ]
        if len(members) != 2:
            continue
        members.sort(key=lambda m: m[1][1])  # by top
        (u_shape, u_box, u_pt), (l_shape, l_box, l_pt) = members

        # Must be a true vertical stack (same column), not two side-by-side boxes.
        x_overlap = min(u_box[0] + u_box[2], l_box[0] + l_box[2]) - max(u_box[0], l_box[0])
        if x_overlap <= 0 or x_overlap < min(u_box[2], l_box[2]) * 0.45:
            continue
        # KPI signature: the large value sits on top of the smaller label.
        if u_pt < l_pt:
            continue

        current_gap = l_box[1] - (u_box[1] + u_box[3])
        value_pt = max(u_pt, l_pt) or 18.0
        target_gap = int(max(5.0, min(14.0, value_pt * 0.25)) * _EMU_PER_PT)
        if current_gap >= target_gap:
            continue

        need = target_gap - current_gap
        room_below = (cb - pad_v) - (l_box[1] + l_box[3])
        room_above = u_box[1] - (ct + pad_v)
        d_down = max(0, min(need, room_below))
        d_up = max(0, min(need - d_down, room_above))
        if d_down + d_up < spacing.min_y_overlap:
            continue  # not enough room to make a meaningful difference

        if d_down > 0:
            try:
                l_shape.top = Emu(l_box[1] + d_down)
                changed = True
            except Exception:
                pass
        if d_up > 0:
            try:
                u_shape.top = Emu(u_box[1] - d_up)
                changed = True
            except Exception:
                pass
    return changed


def _reflow_offslide(slide, slide_w: int, slide_h: int, target_bottom_override: int | None = None) -> bool:
    """
    Lift content that spills past the bottom edge back onto the slide.

    Some templates author a bottom row of cards a little too low, so it runs off
    the slide. The vision QC can't see this (screenshots clip it); only the object
    model shows it. We slide the lowest block of shapes upward into the nearest
    whitespace band — as a RIGID group (every shape at/below a cut moves by the
    same delta) — so each card keeps its internal layout intact. Movement is only
    ever upward, never into the header band, and only by as much as the whitespace
    above allows, so this cannot scramble a slide; worst case it does nothing.
    Footers/page numbers and the header/logo band are left untouched.

    `target_bottom_override` lets the footer-intrusion pass aim higher than the
    slide edge (i.e. at the footer band's top) so cards are lifted clear of the
    footer, not merely back on-slide.
    """
    if slide_w <= 0 or slide_h <= 0:
        return False
    spacing = _Spacing(slide_w=slide_w, slide_h=slide_h)
    bottom_margin = spacing.pad_v
    header_zone_bottom = spacing.header_margin
    footer_zone_top = slide_h - spacing.footer_margin
    min_gap = spacing.gap
    slide_area = max(1, slide_w * slide_h)
    target_bottom = (
        target_bottom_override
        if target_bottom_override is not None
        else slide_h - bottom_margin
    )

    def movable_boxes():
        out = []
        for shape in slide.shapes:
            box = _shape_box(shape)
            if box is None:
                continue
            l, t, w, h = box
            if t >= footer_zone_top:          # footers / page numbers
                continue
            if t < header_zone_bottom:          # header / logo / title band
                continue
            if w * h > slide_area * 0.6:         # full-slide background / panel
                continue
            out.append((shape, box))
        return out

    moved_any = False
    for _ in range(6):
        items = movable_boxes()
        if not items:
            break
        max_bottom = max(b[1] + b[3] for _s, b in items)
        if max_bottom <= target_bottom:
            break  # everything already on-slide
        overflow = max_bottom - target_bottom
        off_top = min(
            (b[1] for _s, b in items if b[1] + b[3] > target_bottom), default=None
        )
        if off_top is None:
            break
        # Largest whitespace band (in 1-D vertical projection) strictly above the
        # off-slide content — that is the room we can slide the block up into.
        running_bottom = None
        best_gap = 0
        best_cut = None
        for _s, b in sorted(items, key=lambda it: it[1][1]):
            t = b[1]
            if running_bottom is not None and running_bottom < t <= off_top:
                gap = t - running_bottom
                if gap > best_gap:
                    best_gap, best_cut = gap, t
            bottom = b[1] + b[3]
            running_bottom = bottom if running_bottom is None else max(running_bottom, bottom)
        if best_cut is None or best_gap <= min_gap:
            break
        lift = min(overflow, best_gap - min_gap)
        block = [(s, b) for s, b in items if b[1] >= best_cut]
        block_top = min(b[1] for _s, b in block)
        if block_top - lift < header_zone_bottom:
            lift = block_top - header_zone_bottom
        if lift <= 0:
            break
        for s, b in block:
            try:
                s.top = Emu(b[1] - lift)
                moved_any = True
            except Exception:
                pass
    return moved_any


def _compress_cluster_x(cluster: list[list[Any]], residual: int, spacing: _Spacing, *, direction: int) -> bool:
    """
    Narrow a horizontal row of shapes by up to `residual` EMU, shrinking the gaps
    between members (never below a minimum). `direction < 0` anchors the leftmost
    member and pulls the rest left; `direction > 0` anchors the rightmost and pulls
    the rest right. Mutates each member's tracked box. Used as a fallback when a
    plain shift can't pull a row fully off a side panel.
    """
    members = sorted(cluster, key=lambda m: m[1][0])
    if len(members) < 2 or residual <= 0:
        return False
    min_gap = max(1, spacing.pad_h // 2)
    slack = []
    for i in range(len(members) - 1):
        a, b = members[i][1], members[i + 1][1]
        slack.append(max(0, (b[0] - (a[0] + a[2])) - min_gap))
    total_slack = sum(slack)
    if total_slack <= 0:
        return False
    take = min(residual, total_slack)
    reductions = [take * s / total_slack for s in slack]
    changed = False
    if direction < 0:
        cum = 0.0
        for i in range(1, len(members)):
            cum += reductions[i - 1]
            s, b = members[i]
            try:
                s.left = Emu(int(b[0] - cum))
                b[0] = int(b[0] - cum)
                changed = True
            except Exception:
                pass
    else:
        cum = 0.0
        for i in range(len(members) - 2, -1, -1):
            cum += reductions[i]
            s, b = members[i]
            try:
                s.left = Emu(int(b[0] + cum))
                b[0] = int(b[0] + cum)
                changed = True
            except Exception:
                pass
    return changed


def _resolve_panel_intrusions(slide, slide_w: int, slide_h: int, spacing: _Spacing) -> bool:
    """
    Pull content that spills onto a side panel back into the content area.

    For each tall side panel we find the row of content straddling its inner edge
    from the outside, then slide that whole row away from the panel (the boxes
    keep their relative layout). If a plain shift can't fully clear it (the row is
    already near the opposite margin), we compress the gaps between the row's boxes
    as a fallback. Vertical layout is never touched.
    """
    all_regions = _keepout_regions(slide, spacing)
    side = [r for r in all_regions if r[0] in ("panel_left", "panel_right")]
    if not side:
        return False
    region_boxes = [r[1] for r in all_regions]
    intruders = _intruder_boxes(slide, spacing, region_boxes)
    tol = spacing.gap
    margin = spacing.pad_h
    header_cut = max(
        (b[1] + b[3] for k, b in all_regions if k == "band_top"),
        default=int(spacing.slide_h * 0.10),
    )
    # Header-strip content (e.g. a wide title) legitimately spans across a
    # same-coloured side panel; never drag it around chasing a phantom overlap.
    body = [(s, b) for (s, b) in intruders if (b[1] + b[3] / 2.0) > header_cut]
    changed = False

    for kind, region in side:
        rl, rt, rw, rh = region
        rr = rl + rw
        straddlers = [(s, b) for (s, b) in body if _intrusion_amount(kind, region, b, tol) > tol]
        if not straddlers:
            continue
        band_top = min(b[1] for _s, b in straddlers)
        band_bottom = max(b[1] + b[3] for _s, b in straddlers)

        cluster: list[list[Any]] = []
        for s, b in body:
            l, t, w, h = b
            v_overlap = min(t + h, band_bottom) - max(t, band_top)
            if v_overlap <= h * 0.5:
                continue
            cx = l + w / 2.0
            if kind == "panel_right" and cx >= rl:
                continue  # this box sits on the panel itself
            if kind == "panel_left" and cx <= rr:
                continue
            cluster.append([s, list(b)])
        if not cluster:
            continue

        if kind == "panel_right":
            max_right = max(b[0] + b[2] for _s, b in cluster)
            overflow = max_right - (rl - margin)
            if overflow <= tol:
                continue
            left_bound = margin
            for k2, r2 in side:
                if k2 == "panel_left":
                    left_bound = max(left_bound, r2[0] + r2[2] + margin)
            cluster_min_left = min(b[0] for _s, b in cluster)
            shift = max(0, min(overflow, cluster_min_left - left_bound))
            if shift > 0:
                for s, b in cluster:
                    try:
                        s.left = Emu(b[0] - shift)
                        b[0] -= shift
                        changed = True
                    except Exception:
                        pass
            residual = overflow - shift
            if residual > tol and _compress_cluster_x(cluster, residual, spacing, direction=-1):
                changed = True
        else:  # panel_left -> push right
            min_left = min(b[0] for _s, b in cluster)
            overflow = (rr + margin) - min_left
            if overflow <= tol:
                continue
            right_bound = slide_w - margin
            for k2, r2 in side:
                if k2 == "panel_right":
                    right_bound = min(right_bound, r2[0] - margin)
            cluster_max_right = max(b[0] + b[2] for _s, b in cluster)
            shift = max(0, min(overflow, right_bound - cluster_max_right))
            if shift > 0:
                for s, b in cluster:
                    try:
                        s.left = Emu(b[0] + shift)
                        b[0] += shift
                        changed = True
                    except Exception:
                        pass
            residual = overflow - shift
            if residual > tol and _compress_cluster_x(cluster, residual, spacing, direction=1):
                changed = True
    return changed


# ── Card clusters (filled card background + the text sitting inside it) ──────
#
# Many templates build a "card" as a filled rectangle/rounded-rectangle with the
# text in SEPARATE text boxes layered on top. The text-only passes above are
# blind to those backgrounds, so three real defects slip through:
#   * two card backgrounds overlapping each other (no text in either),
#   * a stray text box poking into a neighbouring card it doesn't belong to,
#   * a card's own body text rendering past the card's bottom edge.
# Treating "card background + the text geometrically inside it" as ONE movable
# unit lets both detector and fixer act on these. Detection is purely geometric
# (a filled rectangle/rounded-rect that actually holds text) — never by shape
# name — so it generalises to any template. A background only becomes a card if
# it contains text, which automatically excludes decorative ovals/lines.


def _is_card_background(shape) -> bool:
    """A rectangle/rounded-rectangle autoshape or a picture — a plausible card
    backing. Decorative ovals/arcs/diamonds and connector lines are rejected."""
    try:
        st = shape.shape_type
    except Exception:
        return False
    if st == MSO_SHAPE_TYPE.PICTURE:
        return True
    if st != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return False
    try:
        name = str(shape.auto_shape_type)
    except Exception:
        return False
    return "RECTANGLE" in name.upper()


def _collect_card_backgrounds(slide, spacing: _Spacing) -> list[tuple[Any, list[int]]]:
    """Filled card backings: text-less rectangles/pictures within a sane size
    band, excluding the slide background and any structural panel/footer band."""
    slide_area = max(1, spacing.slide_w * spacing.slide_h)
    min_side = spacing.min_container_side
    region_boxes = [r[1] for r in _keepout_regions(slide, spacing)]
    out: list[tuple[Any, list[int]]] = []
    for shape in slide.shapes:
        if _has_visible_text(shape):
            continue
        if not _is_card_background(shape):
            continue
        box = _shape_box(shape)
        if box is None or box in region_boxes:
            continue
        if box[2] < min_side or box[3] < min_side:
            continue
        if box[2] * box[3] > slide_area * 0.6:  # full-slide background
            continue
        out.append((shape, box))
    return out


def _build_clusters(slide, spacing: _Spacing) -> tuple[list[dict[str, Any]], list[tuple[Any, list[int]]]]:
    """
    Group each card background with the text boxes whose centre sits inside it.

    Returns ``(clusters, free_text)`` where each cluster is
    ``{"bg": shape, "box": [l, t, w, h], "members": [(shape, box), ...]}`` and
    ``free_text`` are text boxes not claimed by any card. Smaller cards claim
    their text first so a box lands in the tightest enclosing card.
    """
    cards = _collect_card_backgrounds(slide, spacing)
    text_items: list[tuple[Any, list[int]]] = []
    for shape in slide.shapes:
        if not _has_visible_text(shape):
            continue
        box = _shape_box(shape)
        if box is None:
            continue
        text_items.append((shape, box))

    slack = spacing.slack
    below_tol = spacing.below_tol
    clusters: list[dict[str, Any]] = []
    claimed: set[int] = set()
    for bg, cbox in sorted(cards, key=lambda cb: cb[1][2] * cb[1][3]):
        cl, ct, cw, ch = cbox
        cr, cb = cl + cw, ct + ch
        members: list[tuple[Any, list[int]]] = []
        for shape, tbox in text_items:
            if id(shape) in claimed:
                continue
            tcx = tbox[0] + tbox[2] / 2.0
            tcy = tbox[1] + tbox[3] / 2.0
            if cl - slack <= tcx <= cr + slack and ct - slack <= tcy <= cb + below_tol:
                members.append((shape, tbox))
        if members:
            for shape, _b in members:
                claimed.add(id(shape))
            clusters.append({"bg": bg, "box": list(cbox), "members": members})

    free_text = [(s, b) for (s, b) in text_items if id(s) not in claimed]
    return clusters, free_text


def _enable_word_wrap(shape) -> None:
    try:
        tf = shape.text_frame
        if tf.word_wrap is not True:
            tf.word_wrap = True
    except Exception:
        pass


def _shift_cluster_v(cluster: dict[str, Any], dy: int) -> bool:
    """Move a whole cluster (background + member text) vertically by ``dy`` EMU."""
    if dy == 0:
        return False
    moved = False
    box = cluster["box"]
    try:
        cluster["bg"].top = Emu(int(box[1] + dy))
        box[1] += dy
        moved = True
    except Exception:
        pass
    for shape, tbox in cluster["members"]:
        try:
            shape.top = Emu(int(tbox[1] + dy))
            tbox[1] += dy
            moved = True
        except Exception:
            pass
    return moved


def _resolve_cluster_overlaps(clusters: list[dict[str, Any]], slide_h: int, spacing: _Spacing) -> bool:
    """
    Separate vertically-overlapping card clusters by moving each as a rigid unit.

    The lower cluster is pushed down toward (never past) the footer band; if that
    can't fully clear the overlap, the upper cluster is pulled up toward (never
    into) the header band. Mirrors the bidirectional logic of ``_resolve_container``
    but operates on whole cards, so a card's internal layout is preserved.
    """
    if len(clusters) < 2:
        return False
    footer_zone_top = slide_h - spacing.footer_margin
    header_zone_bottom = spacing.header_margin
    gap = spacing.gap
    changed = False
    for _ in range(6):
        clusters.sort(key=lambda c: c["box"][1])
        moved = False
        n = len(clusters)
        for i in range(n):
            for j in range(i + 1, n):
                A = clusters[i]["box"]
                B = clusters[j]["box"]
                x_overlap = min(A[0] + A[2], B[0] + B[2]) - max(A[0], B[0])
                if x_overlap <= 0 or x_overlap < min(A[2], B[2]) * 0.45:
                    continue
                upper, lower = (clusters[i], clusters[j]) if A[1] <= B[1] else (clusters[j], clusters[i])
                ub, lb = upper["box"], lower["box"]
                y_overlap = (ub[1] + ub[3]) - lb[1]
                if y_overlap <= spacing.min_y_overlap:
                    continue
                need = y_overlap + gap
                room_down = max(0, footer_zone_top - (lb[1] + lb[3]))
                d_down = int(min(need, room_down))
                if d_down > 0 and _shift_cluster_v(lower, d_down):
                    moved = changed = True
                residual = need - d_down
                if residual > spacing.min_y_overlap:
                    room_up = max(0, ub[1] - header_zone_bottom)
                    d_up = int(min(residual, room_up))
                    if d_up > 0 and _shift_cluster_v(upper, -d_up):
                        moved = changed = True
        if not moved:
            break
    return changed


def _resolve_card_text_intrusions(
    clusters: list[dict[str, Any]], free_text: list[tuple[Any, list[int]]], spacing: _Spacing
) -> bool:
    """
    Pull a stray text box back out of a neighbouring card it overlaps from a side.

    The box's centre is outside the card (so it belongs to another column), but an
    edge crosses into the card. We trim/move only that text box — never the card —
    keeping the box's far edge fixed and enabling wrap so the text simply reflows.
    Horizontal only; vertical layout is untouched.
    """
    if not clusters or not free_text:
        return False
    gap = spacing.gap
    changed = False
    for shape, tbox in free_text:
        for cluster in clusters:
            cl, ct, cw, ch = cluster["box"]
            cr, cb = cl + cw, ct + ch
            tl, tt, tw, th = tbox
            tr = tl + tw
            tcx = tl + tw / 2.0
            if min(tt + th, cb) - max(tt, ct) <= th * 0.3:
                continue  # no real vertical overlap with the card
            if min(tr, cr) - max(tl, cl) <= gap:
                continue  # not actually crossing the card edge
            if cl <= tcx <= cr:
                continue  # centre over the card -> not a side intrusion
            floor_w = max(spacing.min_container_side, int(tw * 0.5))
            if tcx < cl:  # poking in from the card's left -> shrink the right edge
                new_w = (cl - gap) - tl
                if floor_w <= new_w < tw:
                    _enable_word_wrap(shape)
                    try:
                        shape.width = Emu(int(new_w))
                        tbox[2] = int(new_w)
                        changed = True
                    except Exception:
                        pass
            else:  # poking in from the right -> push the left edge past the card
                new_left = cr + gap
                new_w = tr - new_left
                if floor_w <= new_w and new_left > tl:
                    _enable_word_wrap(shape)
                    try:
                        shape.left = Emu(int(new_left))
                        shape.width = Emu(int(new_w))
                        tbox[0] = int(new_left)
                        tbox[2] = int(new_w)
                        changed = True
                    except Exception:
                        pass
    return changed


def _normalize_chart_legends(slide) -> bool:
    """
    Give every chart legend its own space instead of overlaying the plot.

    A legend with ``include_in_layout`` true (``c:overlay val="1"``) is drawn ON
    TOP of the plot area, so it visually overshadows the bars/slices. Setting it
    false makes the plot shrink to leave room for the legend. Purely beneficial
    and idempotent; the verify gate is the backstop if a deck ever disagrees.
    """
    changed = False
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        try:
            chart = shape.chart
            if chart.has_legend and chart.legend.include_in_layout is not False:
                chart.legend.include_in_layout = False
                changed = True
        except Exception:
            pass
    return changed


def _shrink_text_to_height(shape, max_h: int, *, min_pt: int = 9) -> bool:
    """
    Shrink a text box's fonts (1pt at a time) until the rendered text actually
    fits `max_h`, measured with the same `_effective_text_height` the detector
    uses — not a single fixed step. Tightens paragraph spacing first, sets the box
    to `max_h`, and enables SHRINK-to-fit as a renderer-side safety. Stops at a
    readable floor (`min_pt`) even if the text still slightly overflows (better a
    small overflow than 6pt text); the verify gate is the final backstop.
    """
    if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
        return False
    if max_h <= 0:
        return False
    changed = False
    tf = shape.text_frame
    try:
        if tf.auto_size != MSO_AUTO_SIZE.SHRINK_TEXT_ON_OVERFLOW:
            tf.auto_size = MSO_AUTO_SIZE.SHRINK_TEXT_ON_OVERFLOW
            changed = True
        if tf.word_wrap is not True:
            tf.word_wrap = True
            changed = True
    except Exception:
        pass
    for p in tf.paragraphs:
        try:
            if p.space_before != Pt(0):
                p.space_before = Pt(0)
                changed = True
            if p.space_after != Pt(0):
                p.space_after = Pt(0)
                changed = True
            if p.line_spacing and isinstance(p.line_spacing, (int, float)) and p.line_spacing > 1.0:
                p.line_spacing = 1.0
                changed = True
        except Exception:
            pass

    width = int(shape.width)
    try:
        shape.height = Emu(int(max_h))
        changed = True
    except Exception:
        pass

    # Iteratively step fonts down until the measured text height fits.
    for _ in range(16):
        if _effective_text_height(shape, width, int(max_h)) <= max_h:
            break
        reduced = False
        for p in tf.paragraphs:
            runs = list(p.runs)
            if runs:
                for run in runs:
                    cur = _cur_run_pt(run, p)
                    new_pt = max(float(min_pt), cur - 1.0)
                    if new_pt < cur:
                        try:
                            run.font.size = Pt(new_pt)
                            reduced = True
                            changed = True
                        except Exception:
                            pass
            else:
                cur = _paragraph_font_pt(p)
                new_pt = max(float(min_pt), cur - 1.0)
                if new_pt < cur:
                    try:
                        p.font.size = Pt(new_pt)
                        reduced = True
                        changed = True
                    except Exception:
                        pass
        if not reduced:
            break  # already at the floor
    return changed


def _cur_run_pt(run, paragraph) -> float:
    """Current font size (pt) of a run, falling back to the paragraph's size."""
    try:
        if run.font.size and run.font.size.pt:
            return float(run.font.size.pt)
    except Exception:
        pass
    return _paragraph_font_pt(paragraph)


def _resolve_footer_intrusion(slide, slide_w: int, slide_h: int, spacing: _Spacing) -> bool:
    """
    Clear content that bleeds into the bottom footer band.

    Two escalating, non-destructive-first strategies:
      1. LIFT — slide the bottom block up into whitespace above (reusing the
         off-slide reflow, but aimed at the footer line, not the slide edge).
      2. SHRINK — for whatever still pokes below the footer line, shrink the card
         height and step its text fonts down so it fits between its top and the
         footer. This is the bounded last-resort resize.
    Footer text itself (center inside the band) is never touched.
    """
    bands = [r for r in _keepout_regions(slide, spacing) if r[0] == "band_bottom"]
    if not bands:
        return False
    footer_line = min(r[1][1] for r in bands)  # highest band top = the footer line
    gap = spacing.gap

    region_boxes = [r[1] for r in _keepout_regions(slide, spacing)]
    intruders = _intruder_boxes(slide, spacing, region_boxes)

    def eff_bottom(shape, box):
        # Use rendered text height for text boxes: an over-stuffed card body whose
        # declared box fits above the footer can still RENDER past it.
        if _has_visible_text(shape):
            return box[1] + _effective_text_height(shape, box[2], box[3])
        return box[1] + box[3]

    straddling = [
        (s, b) for (s, b) in intruders
        if eff_bottom(s, b) > footer_line + gap and (b[1] + (eff_bottom(s, b) - b[1]) / 2.0) < footer_line
    ]
    if not straddling:
        return False

    changed = False
    # 1) Lift the bottom block toward the footer line.
    if _reflow_offslide(slide, slide_w, slide_h, target_bottom_override=footer_line - gap):
        changed = True

    # 2) Shrink whatever still pokes below the footer line (by rendered height).
    for shape in slide.shapes:
        box = _shape_box(shape)
        if box is None or box in region_boxes:
            continue  # never resize a structural panel/band itself
        l, t, w, h = box
        if h > slide_h * 0.5:
            continue  # full-height sidebar/background — legitimately reaches bottom
        e_bottom = eff_bottom(shape, box)
        if e_bottom <= footer_line - gap:
            continue
        if (t + (e_bottom - t) / 2.0) >= footer_line:
            continue  # belongs to the footer band
        new_h = footer_line - gap - t
        if new_h <= 0:
            continue
        if _has_visible_text(shape):
            if _shrink_text_to_height(shape, new_h):
                changed = True
        else:  # a card background — just shorten it
            if new_h >= spacing.min_container_side and new_h < h:
                try:
                    shape.height = Emu(int(new_h))
                    changed = True
                except Exception:
                    pass

    # Pull any in-card text back inside the (now shorter) cards.
    if _fit_text_in_containers(slide, _slide_containers(slide, spacing), spacing):
        changed = True
    return changed


def _fit_card_member_text(clusters: list[dict[str, Any]], spacing: _Spacing) -> bool:
    """
    Shrink a card's own text when it renders past the card's bottom edge.

    Uses the same ``_effective_text_height`` the detector uses, then steps the
    fonts down (to a readable floor) until the text fits between its top and the
    card bottom. Bounded and renderer-safe; resolves "text overflowing its box"
    (e.g. an over-stuffed summary card) without disturbing neighbouring content.
    """
    changed = False
    pad_v = spacing.pad_v
    for cluster in clusters:
        cl, ct, cw, ch = cluster["box"]
        cb = ct + ch
        for shape, tbox in cluster["members"]:
            eff_bottom = tbox[1] + _effective_text_height(shape, tbox[2], tbox[3])
            if eff_bottom <= cb - pad_v:
                continue
            max_h = (cb - pad_v) - tbox[1]
            if max_h <= spacing.min_y_overlap:
                continue
            if _shrink_text_to_height(shape, int(max_h)):
                tbox[3] = int(max_h)
                changed = True
    return changed


def count_card_collisions(slide, spacing: _Spacing) -> int:
    """
    Deterministic count of card-level defects on a slide (engine-agnostic signal):
      * overlapping card backgrounds,
      * a card's body text rendering past the card bottom,
      * a stray text box poking into a foreign card from a side.

    Mirrors exactly what ``_resolve_cluster_overlaps`` / ``_fit_card_member_text``
    / ``_resolve_card_text_intrusions`` act on, so detection and repair stay in
    lock-step. Degrades to 0 on any error.
    """
    try:
        clusters, free_text = _build_clusters(slide, spacing)
    except Exception:
        return 0
    count = 0
    n = len(clusters)
    for i in range(n):
        A = clusters[i]["box"]
        for j in range(i + 1, n):
            B = clusters[j]["box"]
            x_overlap = min(A[0] + A[2], B[0] + B[2]) - max(A[0], B[0])
            if x_overlap <= 0 or x_overlap < min(A[2], B[2]) * 0.45:
                continue
            if min(A[1] + A[3], B[1] + B[3]) - max(A[1], B[1]) > spacing.min_y_overlap:
                count += 1

    pad_v = spacing.pad_v
    for cluster in clusters:
        cl, ct, cw, ch = cluster["box"]
        cb = ct + ch
        for shape, tbox in cluster["members"]:
            eff_bottom = tbox[1] + _effective_text_height(shape, tbox[2], tbox[3])
            if eff_bottom > cb - pad_v + spacing.min_y_overlap:
                count += 1
                break

    gap = spacing.gap
    for shape, tbox in free_text:
        tl, tt, tw, th = tbox
        tr = tl + tw
        tcx = tl + tw / 2.0
        for cluster in clusters:
            cl, ct, cw, ch = cluster["box"]
            cr, cb = cl + cw, ct + ch
            if min(tt + th, cb) - max(tt, ct) <= th * 0.3:
                continue
            if min(tr, cr) - max(tl, cl) <= gap:
                continue
            if cl <= tcx <= cr:
                continue
            count += 1
            break
    return count


def _resolve_text_overlaps(slide, slide_w: int, slide_h: int, *, aggressive: bool) -> bool:
    """
    Geometric (non-destructive) de-overlap pass: move overlapping text boxes
    apart instead of shrinking fonts or trimming content. Runs on the slide's
    top-level shapes and recurses into each group (in the group's own child
    coordinate space) so grouped layouts (e.g. section dividers) are handled too.
    """
    spacing = _Spacing(slide_w=slide_w, slide_h=slide_h)
    changed = False
    # Resolve over-stuffed cards bleeding into the footer FIRST. Otherwise the
    # de-overlap pass below sees a card body overlapping the footer text, pulls the
    # body up into its own title, then "fixes" that by shoving the title down into
    # the footer — ejecting it from the card. Clearing the footer up front (lift +
    # shrink-to-fit) removes that trigger entirely.
    if _resolve_footer_intrusion(slide, slide_w, slide_h, spacing):
        changed = True

    # Containers may have shifted/shortened above, so (re)derive them now.
    containers = _slide_containers(slide, spacing)
    if _resolve_container(slide.shapes, slide_h, spacing, aggressive=aggressive, containers=containers):
        changed = True
    for shape in slide.shapes:
        try:
            is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if not is_group:
            continue
        try:
            child_bottom = _group_child_bottom(shape, slide_h)
            if _resolve_container(shape.shapes, child_bottom, spacing, aggressive=aggressive):
                changed = True
        except Exception:
            pass
    # Tidy text within cards/tiles (pull below-card boxes back in, pad left text).
    if _fit_text_in_containers(slide, containers, spacing):
        changed = True
    # Add breathing room between a KPI value and its caption (intra-card padding).
    if _balance_kpi_cards(slide, containers, spacing):
        changed = True
    # Push content off any side panel it spills onto (horizontal de-collision).
    if _resolve_panel_intrusions(slide, slide_w, slide_h, spacing):
        changed = True
    # Card-aware collisions: a filled card background + the text inside it are one
    # movable unit. Catches/repairs defects the text-only passes miss — overlapping
    # card backgrounds, a card's body text overflowing the card, and a stray text
    # box poking into a neighbouring card. Cards are detected geometrically (a
    # filled rectangle/rounded-rect holding text), never by name, so this is
    # template-agnostic, and every move stays inside the slide's safe zones.
    clusters, free_text = _build_clusters(slide, spacing)
    if _fit_card_member_text(clusters, spacing):
        changed = True
    if _resolve_cluster_overlaps(clusters, slide_h, spacing):
        changed = True
    if _resolve_card_text_intrusions(clusters, free_text, spacing):
        changed = True
    return changed


def clean_template_from_qc(
    pptx_bytes: bytes,
    qc_report: dict[str, Any],
    *,
    repair_level: int = 1,
) -> tuple[bytes, dict[str, Any]]:
    """
    Apply deterministic cleanup directly on an existing PPTX.
    """
    pres = Presentation(io.BytesIO(pptx_bytes))
    slides = qc_report.get("slides") if isinstance(qc_report.get("slides"), list) else []
    aggressive = repair_level >= 2
    slide_w = int(pres.slide_width or 0)
    slide_h = int(pres.slide_height or 0)

    cleaned_entries: list[dict[str, Any]] = []
    attempted_slides = 0
    for idx, slide in enumerate(pres.slides):
        slide_number = idx + 1
        failed = any(
            isinstance(s, dict)
            and int(s.get("slide_number", -1)) == slide_number
            and str(s.get("status", "")).lower() == "fail"
            for s in slides
        )

        actions: list[str] = []

        # Geometric de-overlap runs on EVERY slide. It is non-destructive (only
        # separates text boxes that genuinely overlap once rendered) and catches
        # subtle overlaps the vision QC may not flag.
        if slide_h > 0 and _resolve_text_overlaps(slide, slide_w, slide_h, aggressive=aggressive):
            actions.append("shapes_repositioned")

        # Rescue any content that spills off the bottom of the slide (clipping).
        # Non-destructive: slides off-slide blocks up into whitespace above.
        if slide_h > 0 and _reflow_offslide(slide, slide_w, slide_h):
            actions.append("reflow_offslide")

        # Give every chart legend its own space so it never overlays the plot.
        # Purely beneficial + idempotent, so it runs on every slide.
        if _normalize_chart_legends(slide):
            actions.append("chart_legend_normalized")

        issue_types: list[str] = []
        if failed:
            attempted_slides += 1
            issue_types = _issue_types_for_slide(qc_report, slide_number)
            has_text_issue = any(t in ("overflow", "overlap", "clipping", "alignment") for t in issue_types)
            has_table_issue = any(t in ("table_bounds", "clipping", "overflow") for t in issue_types)
            has_chart_issue = any(t in ("chart_labels", "clipping", "overlap") for t in issue_types)
            if not issue_types:
                has_text_issue = has_table_issue = has_chart_issue = True

            for shape in list(slide.shapes):
                if has_table_issue and _clean_table_shape(shape, aggressive=aggressive):
                    actions.append("table_adjusted")
                if has_chart_issue and _clean_chart_shape(shape, aggressive=aggressive):
                    actions.append("chart_adjusted")
                if has_text_issue and _reduce_text_shape(shape, aggressive=aggressive):
                    actions.append("text_adjusted")

        if actions:
            cleaned_entries.append(
                {
                    "slide_number": slide_number,
                    "issue_types": issue_types,
                    "actions": sorted(set(actions)),
                    "applied": True,
                }
            )

    out = io.BytesIO()
    pres.save(out)
    report = {
        "applied": any(e.get("applied") for e in cleaned_entries),
        "repair_level": repair_level,
        "attempted_slide_count": attempted_slides,
        "cleaned_slide_count": sum(1 for e in cleaned_entries if e.get("applied")),
        "slides": cleaned_entries,
    }
    return out.getvalue(), report
