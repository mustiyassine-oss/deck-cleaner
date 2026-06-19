"""LLM-assisted geometric fix fallback for unresolved QC slides."""
from __future__ import annotations

import io
import json
import logging
from pathlib import Path
from typing import Any

from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from qc import convert_pdf_to_images, convert_pptx_to_pdf, write_pptx_bytes

log = logging.getLogger("deckgen")


def _extract_json(raw: str) -> str:
    s = (raw or "").strip()
    for prefix in ("```json", "```JSON", "```"):
        if s.startswith(prefix):
            s = s[len(prefix) :].lstrip()
            break
    if s.endswith("```"):
        s = s[:-3].rstrip()
    return s


def _failed_slide_numbers(qc_report: dict[str, Any]) -> set[int]:
    slides = qc_report.get("slides") if isinstance(qc_report.get("slides"), list) else []
    out: set[int] = set()
    for s in slides:
        if not isinstance(s, dict):
            continue
        if str(s.get("status", "")).lower() != "fail":
            continue
        try:
            n = int(s.get("slide_number"))
        except (TypeError, ValueError):
            continue
        if n > 0:
            out.add(n)
    return out


def _slide_issues(qc_report: dict[str, Any], slide_number: int) -> list[dict[str, str]]:
    slides = qc_report.get("slides") if isinstance(qc_report.get("slides"), list) else []
    for s in slides:
        if not isinstance(s, dict):
            continue
        try:
            if int(s.get("slide_number")) != slide_number:
                continue
        except (TypeError, ValueError):
            continue
        issues = s.get("issues") if isinstance(s.get("issues"), list) else []
        out: list[dict[str, str]] = []
        for i in issues:
            if not isinstance(i, dict):
                continue
            out.append(
                {
                    "type": str(i.get("type", "")).strip(),
                    "description": str(i.get("description", "")).strip(),
                    "recommended_fix": str(i.get("recommended_fix", "")).strip(),
                }
            )
        return out
    return []


def _shape_inventory(slide, *, slide_w: int, slide_h: int) -> list[dict[str, Any]]:
    footer_zone_top = int(slide_h * 0.92)
    inv: list[dict[str, Any]] = []
    for shape in slide.shapes:
        try:
            sid = int(shape.shape_id)
            left = int(shape.left)
            top = int(shape.top)
            width = int(shape.width)
            height = int(shape.height)
        except (TypeError, ValueError):
            continue
        if width <= 0 or height <= 0:
            continue
        has_text = bool(getattr(shape, "has_text_frame", False) and shape.has_text_frame)
        preview = ""
        alignment = "unknown"
        if has_text:
            try:
                preview = (shape.text_frame.text or "").strip().replace("\n", " ")[:80]
                paragraphs = shape.text_frame.paragraphs or []
                if paragraphs:
                    a = paragraphs[0].alignment
                    if a == PP_ALIGN.CENTER:
                        alignment = "center"
                    elif a == PP_ALIGN.RIGHT:
                        alignment = "right"
                    elif a == PP_ALIGN.LEFT or a is None:
                        alignment = "left"
            except Exception:
                pass
        inv.append(
            {
                "shape_id": sid,
                "name": str(getattr(shape, "name", "")),
                "shape_type": str(getattr(shape, "shape_type", "")),
                "has_text": has_text,
                "text_preview": preview,
                "alignment": alignment,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "in_footer_zone": top >= footer_zone_top,
                "is_group": bool(getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP),
            }
        )
    return inv


def _plan_for_slide(
    client: OpenAI,
    *,
    model: str,
    slide_number: int,
    image_path: Path,
    issues: list[dict[str, str]],
    inventory: list[dict[str, Any]],
    max_dx: int,
    max_dy: int,
) -> dict[str, Any]:
    # Python <-> API-safe base64 without binascii helper imports.
    import base64

    data_url = f"data:image/png;base64,{base64.b64encode(image_path.read_bytes()).decode('ascii')}"
    system = (
        "You generate geometric PPTX fix plans. Return ONLY valid JSON. "
        "Use only move operations for existing shapes. "
        "Do not invent shape ids. Do not output markdown."
    )
    user_text = json.dumps(
        {
            "slide_number": slide_number,
            "issues": issues,
            "constraints": {
                "allowed_op": "move",
                "max_abs_dx_emu": max_dx,
                "max_abs_dy_emu": max_dy,
                "prefer_small_nudges": True,
                "do_not_touch_footer_or_page_numbers": True,
                "only_move_shapes_with_has_text_true": True,
            },
            "target_schema": {
                "slide_number": "int",
                "ops": [
                    {
                        "op": "move",
                        "shape_id": "int",
                        "dx_emu": "int",
                        "dy_emu": "int",
                        "reason": "short string",
                    }
                ],
            },
            "shape_inventory": inventory,
        }
    )
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": user_text},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            },
        ],
        response_format={"type": "json_object"},
        max_tokens=1500,
        temperature=0.1,
    )
    raw = (resp.choices[0].message.content or "").strip()
    try:
        return json.loads(_extract_json(raw))
    except json.JSONDecodeError:
        return {"slide_number": slide_number, "ops": [], "parse_error": True, "raw": raw[:500]}


def apply_llm_fallback_fixes(
    pptx_bytes: bytes,
    qc_report: dict[str, Any],
    *,
    openai_api_key: str,
    model: str,
    work_dir: str | Path,
    openai_base_url: str | None = None,
    max_slides: int = 8,
) -> tuple[bytes, dict[str, Any]]:
    """
    Fallback pass:
      1) ask LLM for safe move plans on remaining failed slides
      2) validate/clamp operations with strict guardrails
      3) apply resulting moves to PPTX
    """
    failed = sorted(_failed_slide_numbers(qc_report))
    if not failed:
        return pptx_bytes, {
            "applied": False,
            "reason": "no_failed_slides",
            "planned_slide_count": 0,
            "touched_slide_numbers": [],
            "plans": [],
            "rejected_ops": [],
        }

    failed = failed[: max(1, int(max_slides))]
    wd = Path(work_dir)
    wd.mkdir(parents=True, exist_ok=True)
    pptx_path = wd / "llm_fix_input.pptx"
    pdf_dir = wd / "pdf"
    img_dir = wd / "slides_png"

    try:
        write_pptx_bytes(pptx_bytes, pptx_path)
        pdf_path = convert_pptx_to_pdf(pptx_path, pdf_dir)
        images = convert_pdf_to_images(pdf_path, img_dir, dpi=220)
    except Exception as e:
        log.warning("[LLM Fix] Screenshot prep failed: %s", e)
        return pptx_bytes, {
            "applied": False,
            "reason": f"screenshot_prep_failed: {e}",
            "planned_slide_count": 0,
            "touched_slide_numbers": [],
            "plans": [],
            "rejected_ops": [],
        }

    client_kw: dict[str, Any] = {"api_key": openai_api_key, "timeout": 600.0}
    if openai_base_url:
        client_kw["base_url"] = openai_base_url
    client = OpenAI(**client_kw)

    pres = Presentation(io.BytesIO(pptx_bytes))
    slide_w = int(pres.slide_width or 0)
    slide_h = int(pres.slide_height or 0)
    max_dx = int(slide_w * 0.08)
    max_dy = int(slide_h * 0.08)
    max_cum_x = int(slide_w * 0.12)
    max_cum_y = int(slide_h * 0.12)
    footer_zone_top = int(slide_h * 0.92)

    plans: list[dict[str, Any]] = []
    rejected_ops: list[dict[str, Any]] = []
    touched_slides: set[int] = set()
    cumulative: dict[tuple[int, int], tuple[int, int]] = {}
    changed = False

    for slide_no in failed:
        idx = slide_no - 1
        if idx < 0 or idx >= len(pres.slides):
            continue
        if idx >= len(images):
            continue
        slide = pres.slides[idx]
        inv = _shape_inventory(slide, slide_w=slide_w, slide_h=slide_h)
        issues = _slide_issues(qc_report, slide_no)
        plan = _plan_for_slide(
            client,
            model=model,
            slide_number=slide_no,
            image_path=images[idx],
            issues=issues,
            inventory=inv,
            max_dx=max_dx,
            max_dy=max_dy,
        )
        plans.append(plan)
        ops = plan.get("ops") if isinstance(plan.get("ops"), list) else []
        by_id: dict[int, Any] = {}
        for shp in slide.shapes:
            try:
                by_id[int(shp.shape_id)] = shp
            except (TypeError, ValueError):
                continue

        for op in ops:
            if not isinstance(op, dict):
                continue
            if str(op.get("op", "")).lower() != "move":
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "unsupported_op"})
                continue
            try:
                sid = int(op.get("shape_id"))
                dx = int(op.get("dx_emu", 0))
                dy = int(op.get("dy_emu", 0))
            except (TypeError, ValueError):
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "invalid_fields"})
                continue
            shp = by_id.get(sid)
            if shp is None:
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "unknown_shape"})
                continue
            if not (getattr(shp, "has_text_frame", False) and shp.has_text_frame):
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "shape_has_no_text"})
                continue
            try:
                left = int(shp.left)
                top = int(shp.top)
                width = int(shp.width)
                height = int(shp.height)
            except (TypeError, ValueError):
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "non_positionable"})
                continue
            if top >= footer_zone_top:
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "footer_zone_protected"})
                continue
            if abs(dx) > max_dx or abs(dy) > max_dy:
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "delta_exceeds_guardrail"})
                continue
            key = (slide_no, sid)
            prev_x, prev_y = cumulative.get(key, (0, 0))
            if abs(prev_x + dx) > max_cum_x or abs(prev_y + dy) > max_cum_y:
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "cumulative_exceeds_guardrail"})
                continue

            new_left = left + dx
            new_top = top + dy
            new_left = max(0, min(new_left, max(0, slide_w - width)))
            new_top = max(0, min(new_top, max(0, slide_h - height)))
            # Re-check footer protection after clamp.
            if new_top >= footer_zone_top:
                new_top = min(top, max(0, footer_zone_top - height - int(slide_h * 0.01)))

            if new_left == left and new_top == top:
                continue
            try:
                shp.left = new_left
                shp.top = new_top
                cumulative[key] = (prev_x + (new_left - left), prev_y + (new_top - top))
                changed = True
                touched_slides.add(slide_no)
            except Exception:
                rejected_ops.append({"slide": slide_no, "op": op, "reason": "apply_failed"})

    if not changed:
        return pptx_bytes, {
            "applied": False,
            "reason": "no_safe_ops_applied",
            "planned_slide_count": len(plans),
            "touched_slide_numbers": [],
            "plans": plans,
            "rejected_ops": rejected_ops,
        }

    out = io.BytesIO()
    pres.save(out)
    return out.getvalue(), {
        "applied": True,
        "reason": "applied_safe_ops",
        "planned_slide_count": len(plans),
        "touched_slide_numbers": sorted(touched_slides),
        "changed_shape_count": len(cumulative),
        "plans": plans,
        "rejected_ops": rejected_ops,
    }

