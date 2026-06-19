"""Deterministic layout QC — object-model geometry (primary) + PDF geometry (complement).

Two deterministic signal sources are fused:

* **Object model** (`python-pptx`): the *declared* geometry of every shape — exact
  positions/sizes even when a shape sits off-slide or boxes overlap. This is the
  primary detector for shape-vs-shape overlap, off-slide / below-edge tables, and
  tiny declared fonts. It needs no rendering and no external tools.

* **PDF geometry** (PyMuPDF, optional): the *rendered* layout after LibreOffice
  has applied real wrapping/autofit. Complements the object model for things only
  visible once rendered — chart-label collisions and rendered text overlaps. If
  PyMuPDF or LibreOffice is unavailable, this pass is simply skipped and the
  object-model result still stands.

The overlap criterion deliberately mirrors the fixer in ``template_cleaner.py``
(same horizontal-overlap ratio + effective-text-height test), so a slide is
flagged exactly when the deterministic fixer would act on it — detection and
repair stay consistent.

Emits the same report schema as ``qc.run_visual_qc`` so it is a drop-in detector.
No OpenAI key required.
"""
from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from qc import (
    _resolve_soffice,
    _structural_text_issues,
    convert_pptx_to_pdf,
    write_pptx_bytes,
)
from template_cleaner import (
    _Spacing,
    _count_keepout_intrusions,
    _effective_text_height,
    _has_visible_text,
    _shape_box,
    count_card_collisions,
)

log = logging.getLogger("cleaner")

_SEVERITY_RANK = {"none": 0, "low": 1, "medium": 2, "high": 3}
_MIN_X_OVERLAP_RATIO = 0.45


def _skipped(reason: str) -> dict[str, Any]:
    return {
        "overall_status": "skipped",
        "skipped_reason": reason,
        "failed_slide_count": 0,
        "total_slide_count": 0,
        "slides": [],
        "engine": "geometry",
    }


def _severity_for(issues: list[dict[str, str]]) -> str:
    if not issues:
        return "none"
    types = {i.get("type") for i in issues}
    if {"overlap", "clipping", "table_bounds", "panel_overlap"} & types:
        return "high"
    return "medium"


def _merge_issues(*groups: list[dict[str, str]]) -> list[dict[str, str]]:
    """Concatenate issue lists, keeping the first occurrence of each type."""
    out: list[dict[str, str]] = []
    seen: set[str] = set()
    for group in groups:
        for issue in group:
            t = str(issue.get("type", ""))
            if not t or t in seen:
                continue
            seen.add(t)
            out.append(issue)
    return out


# ── Object-model detection (primary) ────────────────────────────────────────

def _chart_legend_overlay_count(slide) -> int:
    """Charts whose legend is laid over the plot area (``include_in_layout`` true).

    Such a legend is drawn ON TOP of the bars/slices, so it visually overshadows
    the chart. The vision engine can sometimes flag this, but the object model
    sees it exactly and for free — and the existing chart fixer already resolves
    it once the slide is flagged."""
    n = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        try:
            chart = shape.chart
            if chart.has_legend and chart.legend.include_in_layout is True:
                n += 1
        except Exception:
            pass
    return n



def _collect_text_items(shapes) -> list[tuple[list[int], int]]:
    """[(box, effective_height)] for visible-text shapes (groups handled separately)."""
    items: list[tuple[list[int], int]] = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
        except Exception:
            pass
        if not _has_visible_text(shape):
            continue
        box = _shape_box(shape)
        if box is None:
            continue
        eff = _effective_text_height(shape, box[2], box[3])
        items.append((box, eff))
    return items


def _count_overlaps(items: list[tuple[list[int], int]], spacing: _Spacing) -> int:
    """Count text-box pairs that overlap (same criterion the fixer uses to act)."""
    n = len(items)
    count = 0
    for i in range(n):
        a_box, a_eff = items[i]
        al, at, aw, _ah = a_box
        ar = al + aw
        for j in range(i + 1, n):
            b_box, b_eff = items[j]
            bl, bt, bw, _bh = b_box
            br = bl + bw
            x_overlap = min(ar, br) - max(al, bl)
            if x_overlap <= 0:
                continue
            if x_overlap < min(aw, bw) * _MIN_X_OVERLAP_RATIO:
                continue
            y_overlap = min(at + a_eff, bt + b_eff) - max(at, bt)
            if y_overlap <= spacing.min_y_overlap:
                continue
            count += 1
    return count


def overlap_counts(pptx_bytes: bytes) -> dict[int, int]:
    """
    Per-slide count of overlapping text-box pairs (object model, top-level + groups).

    Used as an engine-agnostic regression check: the vision QC cannot *see*
    geometric overlap, so after a fix pass we compare these counts before/after to
    catch overlaps a fix may have introduced (and roll those slides back).
    """
    out: dict[int, int] = {}
    try:
        pres = Presentation(io.BytesIO(pptx_bytes))
    except Exception:
        return out
    sw = int(pres.slide_width or 0)
    sh = int(pres.slide_height or 0)
    if not (sw and sh):
        return out
    spacing = _Spacing(slide_w=sw, slide_h=sh)
    for idx, slide in enumerate(pres.slides):
        count = _count_overlaps(_collect_text_items(slide.shapes), spacing)
        for shape in slide.shapes:
            try:
                is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
            except Exception:
                is_group = False
            if is_group:
                count += _count_overlaps(_collect_text_items(shape.shapes), spacing)
        out[idx + 1] = count
    return out


def _collect_picture_boxes(shapes) -> list[list[int]]:
    """Boxes of picture shapes (e.g. the logo) — used as collision obstacles."""
    boxes: list[list[int]] = []
    for shape in shapes:
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        box = _shape_box(shape)
        if box is not None:
            boxes.append(box)
    return boxes


def _count_text_picture_overlaps(
    items: list[tuple[list[int], int]], pic_boxes: list[list[int]], spacing: _Spacing
) -> int:
    """Count text boxes that overlap a picture (same criterion as text-text)."""
    count = 0
    for a_box, a_eff in items:
        al, at, aw, _ah = a_box
        ar = al + aw
        for pl, pt, pw, ph in pic_boxes:
            pr = pl + pw
            x_overlap = min(ar, pr) - max(al, pl)
            if x_overlap <= 0:
                continue
            if x_overlap < min(aw, pw) * _MIN_X_OVERLAP_RATIO:
                continue
            y_overlap = min(at + a_eff, pt + ph) - max(at, pt)
            if y_overlap <= spacing.min_y_overlap:
                continue
            count += 1
    return count


def picture_overlap_counts(pptx_bytes: bytes) -> dict[int, int]:
    """
    Per-slide count of text boxes overlapping a picture (e.g. text slid under the
    logo). The vision QC and the text-text overlap check are both blind to
    text-over-image collisions, so the verify gate uses this before/after to catch
    fixes that push text into the logo/header artwork and roll them back.
    """
    out: dict[int, int] = {}
    try:
        pres = Presentation(io.BytesIO(pptx_bytes))
    except Exception:
        return out
    sw = int(pres.slide_width or 0)
    sh = int(pres.slide_height or 0)
    if not (sw and sh):
        return out
    spacing = _Spacing(slide_w=sw, slide_h=sh)
    for idx, slide in enumerate(pres.slides):
        pics = _collect_picture_boxes(slide.shapes)
        out[idx + 1] = (
            _count_text_picture_overlaps(_collect_text_items(slide.shapes), pics, spacing)
            if pics
            else 0
        )
    return out


def _off_slide(slide, sw: int, sh: int, tol_w: int, tol_h: int) -> tuple[int, bool]:
    """
    Count top-level text/table shapes whose declared box leaves the slide, and
    whether a table extends past the bottom edge. Decorative non-text shapes are
    ignored (full-bleed banners/images legitimately cross the edge).
    """
    off = 0
    table_below = False
    for shape in slide.shapes:
        is_table = bool(getattr(shape, "has_table", False))
        if not (is_table or _has_visible_text(shape)):
            continue
        box = _shape_box(shape)
        if box is None:
            continue
        l, t, w, h = box
        if l < -tol_w or t < -tol_h or l + w > sw + tol_w or t + h > sh + tol_h:
            off += 1
            if is_table and t + h > sh + tol_h:
                table_below = True
    return off, table_below


def _object_model_issues(pptx_bytes: bytes) -> tuple[dict[int, list[dict[str, str]]], int]:
    """Per-slide deterministic issues from the PPTX object model. Returns (issues, total_slides)."""
    pres = Presentation(io.BytesIO(pptx_bytes))
    sw = int(pres.slide_width or 0)
    sh = int(pres.slide_height or 0)
    spacing = _Spacing(slide_w=sw, slide_h=sh) if sw and sh else None
    tol_w = max(1, int(sw * 0.012))
    tol_h = max(1, int(sh * 0.012))
    structural = _structural_text_issues(pptx_bytes)

    issues_by_slide: dict[int, list[dict[str, str]]] = {}
    total = 0
    for idx, slide in enumerate(pres.slides):
        total += 1
        slide_no = idx + 1
        issues: list[dict[str, str]] = []

        if spacing is not None:
            overlaps = _count_overlaps(_collect_text_items(slide.shapes), spacing)
            for shape in slide.shapes:
                try:
                    is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
                except Exception:
                    is_group = False
                if is_group:
                    overlaps += _count_overlaps(_collect_text_items(shape.shapes), spacing)
            if overlaps:
                issues.append(
                    {
                        "type": "overlap",
                        "description": f"{overlaps} pair(s) of text boxes overlap once rendered.",
                        "recommended_fix": "Separate the overlapping text boxes.",
                    }
                )

            off, table_below = _off_slide(slide, sw, sh, tol_w, tol_h)
            if off:
                if table_below:
                    issues.append(
                        {
                            "type": "table_bounds",
                            "description": "A table extends past the bottom edge of the slide.",
                            "recommended_fix": "Reduce rows/height or move the table up so it fits.",
                        }
                    )
                else:
                    issues.append(
                        {
                            "type": "clipping",
                            "description": f"{off} text/table element(s) extend beyond the slide bounds.",
                            "recommended_fix": "Move or resize the element back inside the slide.",
                        }
                    )

            intrusions = _count_keepout_intrusions(slide, spacing)
            if intrusions:
                issues.append(
                    {
                        "type": "panel_overlap",
                        "description": f"{intrusions} element(s) overlap a side panel or footer/header band.",
                        "recommended_fix": "Move the element off the panel, or shrink it clear of the band.",
                    }
                )

            # Card-level collisions the text-only checks miss: overlapping card
            # backgrounds, a card's text overflowing the card, or a stray text box
            # poking into a neighbouring card. Same units the fixer repairs.
            cards = count_card_collisions(slide, spacing)
            if cards:
                issues.append(
                    {
                        "type": "overlap",
                        "description": (
                            f"{cards} card collision(s): overlapping cards, text overflowing a card, "
                            "or text poking into a neighbouring card."
                        ),
                        "recommended_fix": "Separate the cards or trim/space the text so each card's content stays inside it.",
                    }
                )

        # Chart legend overlaying the plot (independent of slide geometry).
        if _chart_legend_overlay_count(slide):
            issues.append(
                {
                    "type": "chart_labels",
                    "description": "A chart legend is laid over the plot area and overshadows the chart.",
                    "recommended_fix": "Give the legend its own space so it doesn't cover the plot.",
                }
            )

        for si in structural.get(slide_no, []):
            issues.append(si)

        issues_by_slide[slide_no] = issues
    return issues_by_slide, total


def offslide_issues(pptx_bytes: bytes) -> dict[int, list[dict[str, str]]]:
    """
    Always-on off-slide / clipping detection from the object model only (no PDF
    render needed). Returns per-slide (1-based) clipping/table_bounds issues.

    This is the ONE deterministic signal we feed into the vision engine too:
    screenshots clip off-slide content, so the vision QC is structurally blind to
    it. Object-model coordinates are not. Only clipping is surfaced here — overlap
    detection is left to each engine so we don't change vision's overlap behaviour.
    """
    out: dict[int, list[dict[str, str]]] = {}
    try:
        pres = Presentation(io.BytesIO(pptx_bytes))
    except Exception as e:
        log.warning("[OffSlide] Could not open PPTX: %s", e)
        return out
    sw = int(pres.slide_width or 0)
    sh = int(pres.slide_height or 0)
    if not (sw and sh):
        return out
    tol_w = max(1, int(sw * 0.012))
    tol_h = max(1, int(sh * 0.012))
    for idx, slide in enumerate(pres.slides):
        try:
            off, table_below = _off_slide(slide, sw, sh, tol_w, tol_h)
        except Exception:
            continue
        if not off:
            continue
        if table_below:
            out[idx + 1] = [
                {
                    "type": "table_bounds",
                    "description": "A table extends past the bottom edge of the slide.",
                    "recommended_fix": "Reduce rows/height or move the table up so it fits.",
                }
            ]
        else:
            out[idx + 1] = [
                {
                    "type": "clipping",
                    "description": f"{off} text/table element(s) extend beyond the slide bounds.",
                    "recommended_fix": "Move or resize the element back inside the slide.",
                }
            ]
    return out


def keepout_issues(pptx_bytes: bytes) -> dict[int, list[dict[str, str]]]:
    """
    Always-on side-panel / band intrusion detection from the object model.

    Like ``offslide_issues``, this is fed into the vision engine too: the vision QC
    often misses a card spilling onto a dark sidebar or into the footer bar, but
    the object model sees it exactly. Returns per-slide (1-based) ``panel_overlap``
    issues. Degrades to an empty dict on any error.
    """
    out: dict[int, list[dict[str, str]]] = {}
    try:
        pres = Presentation(io.BytesIO(pptx_bytes))
    except Exception as e:
        log.warning("[KeepOut] Could not open PPTX: %s", e)
        return out
    sw = int(pres.slide_width or 0)
    sh = int(pres.slide_height or 0)
    if not (sw and sh):
        return out
    spacing = _Spacing(slide_w=sw, slide_h=sh)
    for idx, slide in enumerate(pres.slides):
        try:
            n = _count_keepout_intrusions(slide, spacing)
        except Exception:
            continue
        if n:
            out[idx + 1] = [
                {
                    "type": "panel_overlap",
                    "description": f"{n} element(s) overlap a side panel or footer/header band.",
                    "recommended_fix": "Move the element off the panel, or shrink it clear of the band.",
                }
            ]
    return out


# ── PDF rendered-geometry detection (complement) ────────────────────────────

def _rect_area(b: tuple[float, float, float, float]) -> float:
    return max(0.0, b[2] - b[0]) * max(0.0, b[3] - b[1])


def _intersection(
    a: tuple[float, float, float, float],
    b: tuple[float, float, float, float],
) -> tuple[float, float]:
    ix = min(a[2], b[2]) - max(a[0], b[0])
    iy = min(a[3], b[3]) - max(a[1], b[1])
    return (max(0.0, ix), max(0.0, iy))


def _collect_lines(page) -> list[dict[str, Any]]:
    lines: list[dict[str, Any]] = []
    try:
        data = page.get_text("dict")
    except Exception:
        return lines
    for block in data.get("blocks", []):
        if block.get("type", 0) != 0:
            continue
        bidx = block.get("number", -1)
        for ln in block.get("lines", []):
            spans = ln.get("spans", []) or []
            text = "".join(s.get("text", "") for s in spans).strip()
            if not text:
                continue
            size = max((float(s.get("size", 0.0)) for s in spans), default=0.0)
            bbox = ln.get("bbox")
            if not bbox or len(bbox) != 4:
                continue
            lines.append(
                {
                    "bbox": (float(bbox[0]), float(bbox[1]), float(bbox[2]), float(bbox[3])),
                    "size": size,
                    "block": bidx,
                }
            )
    return lines


def _pdf_page_issues(page) -> list[dict[str, str]]:
    """Rendered-geometry defects (complement to the object model)."""
    pr = page.rect
    pw, ph = float(pr.width), float(pr.height)
    if pw <= 0 or ph <= 0:
        return []
    min_inter = max(1.5, ph * 0.0025)

    lines = _collect_lines(page)
    issues: list[dict[str, str]] = []

    overlaps = 0
    n = len(lines)
    for i in range(n):
        a = lines[i]["bbox"]
        a_area = _rect_area(a)
        if a_area <= 0:
            continue
        for j in range(i + 1, n):
            if lines[i]["block"] == lines[j]["block"]:
                continue
            b = lines[j]["bbox"]
            iw, ih = _intersection(a, b)
            if iw <= min_inter or ih <= min_inter:
                continue
            smaller = min(a_area, _rect_area(b))
            if smaller > 0 and (iw * ih) / smaller >= 0.18:
                overlaps += 1
    if overlaps:
        issues.append(
            {
                "type": "overlap",
                "description": f"{overlaps} pair(s) of rendered text blocks overlap.",
                "recommended_fix": "Separate the overlapping text.",
            }
        )
    return issues


def _pdf_issues(pptx_bytes: bytes, work_dir: Path) -> tuple[dict[int, list[dict[str, str]]], str]:
    """Run the optional PDF pass. Returns (issues_by_slide, status_note)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return {}, "skipped: PyMuPDF not installed (pip install PyMuPDF)"
    if not _resolve_soffice():
        return {}, "skipped: LibreOffice not found"

    try:
        pptx_path = work_dir / "deck_for_geoqc.pptx"
        write_pptx_bytes(pptx_bytes, pptx_path)
        pdf_path = convert_pptx_to_pdf(pptx_path, work_dir / "pdf")
        doc = fitz.open(str(pdf_path))
    except Exception as e:
        return {}, f"skipped: {e}"

    out: dict[int, list[dict[str, str]]] = {}
    try:
        for i in range(doc.page_count):
            try:
                page_issues = _pdf_page_issues(doc[i])
            except Exception as e:
                log.warning("[GeoQC] PDF slide %d analysis failed: %s", i + 1, e)
                page_issues = []
            if page_issues:
                out[i + 1] = page_issues
    finally:
        doc.close()
    return out, "ran"


# ── Public entry point ──────────────────────────────────────────────────────

def run_geometry_qc(
    pptx_bytes: bytes,
    work_dir: str | Path,
    *,
    slide_numbers: set[int] | list[int] | tuple[int, ...] | None = None,
    dpi: int = 220,  # accepted for signature parity with run_visual_qc; unused
) -> dict[str, Any]:
    """
    Deterministic layout QC (object-model primary + optional PDF complement).
    On a fatal object-model error, returns a structured skip report (does not
    raise). Emits the same schema as ``qc.run_visual_qc``.
    """
    work_dir = Path(work_dir)
    work_dir.mkdir(parents=True, exist_ok=True)
    target = {int(n) for n in (slide_numbers or []) if int(n) > 0}

    try:
        om_issues, total = _object_model_issues(pptx_bytes)
    except Exception as e:
        log.warning("[GeoQC] Object-model analysis failed: %s", e)
        return _skipped(f"object-model analysis failed: {e}")

    pdf_issues, pdf_status = _pdf_issues(pptx_bytes, work_dir)
    log.info("[GeoQC] PDF complement pass: %s", pdf_status)

    slides_out: list[dict[str, Any]] = []
    failed = 0
    for slide_no in range(1, total + 1):
        if target and slide_no not in target:
            continue
        merged = _merge_issues(om_issues.get(slide_no, []), pdf_issues.get(slide_no, []))
        status = "fail" if merged else "pass"
        severity = _severity_for(merged)
        slides_out.append(
            {
                "slide_number": slide_no,
                "status": status,
                "severity": severity,
                "issues": merged,
            }
        )
        log.info(
            "[GeoQC] Slide %d: status=%s severity=%s issues=%d",
            slide_no,
            status,
            severity,
            len(merged),
        )
        if status == "fail":
            failed += 1

    overall = "fail" if failed else "pass"
    log.info(
        "[GeoQC] Final result: overall_status=%s failed_slide_count=%d total_slide_count=%d",
        overall,
        failed,
        total,
    )
    return {
        "overall_status": overall,
        "failed_slide_count": failed,
        "total_slide_count": total,
        "checked_slide_count": len(slides_out),
        "checked_slide_numbers": [s["slide_number"] for s in slides_out],
        "slides": slides_out,
        "engine": "geometry",
        "pdf_pass": pdf_status,
    }
